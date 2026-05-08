"""MeetPod 소개 PPT 생성 (10 slides, 한글, 비개발자 포함)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x14, 0x2A, 0x5C)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"
    return tb


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(1.0), NAVY)
    add_text(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.35),
                 subtitle, size=12, color=LIGHT)


def footer(slide, page):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "MeetPod — 친구들의 약속을 더 쉽게", size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page} / 10", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ── Slide 1: 표지 ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, 0, Inches(3.0), SW, Inches(0.04), ACCENT)
add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.2),
         "MeetPod", size=72, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.2), Inches(12), Inches(0.6),
         "친구들의 약속을 더 쉽게 — 위치 공유 · 채팅 · 스케줄",
         size=22, color=LIGHT)
add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
         "프로젝트 소개  ·  2026.05.02  ·  CPKWorks",
         size=12, color=MUTED)

# ── Slide 2: 한 줄 요약 ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "한 줄 요약", "MeetPod is …")
add_rect(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.8), LIGHT)
add_rect(s, Inches(1.5), Inches(2.5), Inches(0.15), Inches(2.8), ACCENT)
add_text(s, Inches(2.0), Inches(2.9), Inches(9.5), Inches(0.6),
         "친구들끼리 약속을 잡고,", size=28, bold=True, color=DARK)
add_text(s, Inches(2.0), Inches(3.55), Inches(9.5), Inches(0.6),
         "약속 시간 동안 서로의 위치를 안전하게 공유하며,", size=24, color=DARK)
add_text(s, Inches(2.0), Inches(4.2), Inches(9.5), Inches(0.6),
         "그룹·약속별 채팅으로 모임을 더 쉽게 만드는 앱.", size=24, color=DARK)
add_text(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.5),
         "PickPod(부모-자녀 픽업)와 같은 워크스페이스, 다른 사용자층 — 친구 간 평등 관계와 프라이버시 친화적 위치 공유가 핵심.",
         size=14, color=MUTED)
footer(s, 2)

# ── Slide 3: 왜 만드는가 (문제) ─────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "왜 만드는가", "친구 약속에서 반복되는 불편")

problems = [
    ("📍", "어디야?", "약속 시간 다가오면 카톡·전화로 위치 묻기 반복"),
    ("⏰", "언제 와?", "지각 여부·도착 예상 시각이 불투명"),
    ("🍽️", "어디서 만나지?", "장소 공유는 캡처·링크 복붙으로 산만"),
    ("💬", "대화가 흩어짐", "약속별 대화가 단톡방에 섞여 정보 추적 어려움"),
]
x0, y0, cw, ch, gap = Inches(0.8), Inches(1.6), Inches(5.9), Inches(2.4), Inches(0.3)
for i, (icon, t, desc) in enumerate(problems):
    col, row = i % 2, i // 2
    x = x0 + (cw + gap) * col
    y = y0 + (ch + gap) * row
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.12), ch, ACCENT)
    add_text(s, x + Inches(0.4), y + Inches(0.3), Inches(1.0), Inches(0.7),
             icon, size=36)
    add_text(s, x + Inches(1.4), y + Inches(0.4), cw - Inches(1.6), Inches(0.6),
             t, size=22, bold=True, color=DARK)
    add_text(s, x + Inches(1.4), y + Inches(1.1), cw - Inches(1.6), Inches(1.2),
             desc, size=14, color=MUTED)
footer(s, 3)

# ── Slide 4: 해결 방식 ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "MeetPod의 해결 방식", "약속 = 일정 + 위치 + 채팅, 한 곳에서")

solutions = [
    ("🗓️", "약속 중심 UX", "그룹 또는 1회성 약속을 만들면 일정·장소·참여자가 한 번에 정리됨"),
    ("📡", "스마트 위치 공유", "기본 20분 전 자동 시작, 종료 시 자동 OFF — 사용자가 끄는 걸 잊을 일 없음"),
    ("💬", "약속별 채팅방", "텍스트·이미지·장소(Google Maps) 카드를 약속 단위로 정리, 끝나면 아카이브"),
    ("🔒", "프라이버시 친화", "친구 추가는 초대 링크/QR로만, 위치는 약속 참여자에게만, 핑은 24시간 후 자동 삭제"),
]
x0, y0, cw, ch, gap = Inches(0.8), Inches(1.6), Inches(5.9), Inches(2.4), Inches(0.3)
for i, (icon, t, desc) in enumerate(solutions):
    col, row = i % 2, i // 2
    x = x0 + (cw + gap) * col
    y = y0 + (ch + gap) * row
    add_rect(s, x, y, cw, ch, LIGHT)
    add_rect(s, x, y, Inches(0.12), ch, GREEN)
    add_text(s, x + Inches(0.4), y + Inches(0.3), Inches(1.0), Inches(0.7),
             icon, size=36)
    add_text(s, x + Inches(1.4), y + Inches(0.4), cw - Inches(1.6), Inches(0.6),
             t, size=22, bold=True, color=DARK)
    add_text(s, x + Inches(1.4), y + Inches(1.1), cw - Inches(1.6), Inches(1.2),
             desc, size=14, color=MUTED)
footer(s, 4)

# ── Slide 5: 주요 기능 ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "주요 기능 (MVP)", "최소 기능으로 빠르게, 군더더기 없이")

features = [
    ("로그인 / 가입", ["Google · Apple · Kakao 소셜 로그인", "@핸들 1회 설정", "Expo Push 알림 토큰 등록"]),
    ("친구 / 그룹", ["초대 링크 또는 QR로만 친구 추가", "그룹: 방장 + 관리자 위임 가능", "멤버 초대·추방"]),
    ("약속 (스케줄)", ["제목 · 일시 · 장소(지도) · 참여자", "그룹 약속은 멤버 일괄 선택 + 빠질 사람 해제", "사용자별 개인 알림 (예: 30분 전)"]),
    ("위치 공유", ["기본 20분 전 자동 ON, 종료 시 OFF", "약속별 시작 시각 변경 가능 (10/20/30/60분)", "약속 참여자만 지도에서 핀 확인"]),
    ("채팅", ["그룹 상시 채팅방 + 약속별 채팅방", "텍스트 · 이미지 · 장소 카드", "약속 종료 시 채팅방 자동 아카이브"]),
]
x0, y0, cw, ch, gap = Inches(0.5), Inches(1.5), Inches(2.5), Inches(5.0), Inches(0.05)
for i, (title, items) in enumerate(features):
    x = x0 + (cw + gap) * i
    add_rect(s, x, y0, cw, Inches(0.7), NAVY)
    add_text(s, x + Inches(0.15), y0 + Inches(0.18), cw - Inches(0.3), Inches(0.4),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, y0 + Inches(0.7), cw, ch - Inches(0.7), LIGHT)
    add_bullets(s, x + Inches(0.2), y0 + Inches(0.9), cw - Inches(0.4), ch - Inches(1.0),
                items, size=11, color=DARK, line_spacing=1.3)
footer(s, 5)

# ── Slide 6: PickPod와의 차이 ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "PickPod와의 차이", "같은 워크스페이스, 다른 사용자층")

# Table-style comparison
hdr_y = Inches(1.6)
row_h = Inches(0.7)
col_x = [Inches(0.8), Inches(4.5), Inches(8.5)]
col_w = [Inches(3.5), Inches(3.8), Inches(4.8)]

# Header row
add_rect(s, col_x[0], hdr_y, col_w[0], row_h, NAVY)
add_rect(s, col_x[1], hdr_y, col_w[1], row_h, NAVY)
add_rect(s, col_x[2], hdr_y, col_w[2], row_h, NAVY)
add_text(s, col_x[0] + Inches(0.2), hdr_y + Inches(0.18), col_w[0], Inches(0.4),
         "관점", size=14, bold=True, color=WHITE)
add_text(s, col_x[1] + Inches(0.2), hdr_y + Inches(0.18), col_w[1], Inches(0.4),
         "PickPod", size=14, bold=True, color=WHITE)
add_text(s, col_x[2] + Inches(0.2), hdr_y + Inches(0.18), col_w[2], Inches(0.4),
         "MeetPod", size=14, bold=True, color=WHITE)

rows = [
    ("사용자", "부모 ↔ 자녀", "친구 ↔ 친구"),
    ("관계", "보호 / 모니터링", "평등 / 동의 기반"),
    ("위치 공유", "상시(픽업 시간대 중심)", "약속 시간에만 자동"),
    ("채팅", "필수 아님", "그룹·약속별 1급 기능"),
    ("친구 추가", "보호자가 자녀 등록", "초대 링크 / QR"),
    ("프라이버시", "보호자 권한 우선", "참여자 동의 우선"),
]
for i, (k, a, b) in enumerate(rows):
    y = hdr_y + row_h + row_h * i
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, col_x[0], y, col_w[0], row_h, bg)
    add_rect(s, col_x[1], y, col_w[1], row_h, bg)
    add_rect(s, col_x[2], y, col_w[2], row_h, bg)
    add_text(s, col_x[0] + Inches(0.2), y + Inches(0.2), col_w[0], Inches(0.4),
             k, size=13, bold=True, color=DARK)
    add_text(s, col_x[1] + Inches(0.2), y + Inches(0.2), col_w[1], Inches(0.4),
             a, size=13, color=DARK)
    add_text(s, col_x[2] + Inches(0.2), y + Inches(0.2), col_w[2], Inches(0.4),
             b, size=13, color=DARK)

footer(s, 6)

# ── Slide 7: 기술 스택 / 아키텍처 ───────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "기술 스택 & 아키텍처", "PickPod와 동일 컨벤션으로 빠른 개발")

# Stack boxes
stacks = [
    ("Mobile", "Expo + React Native\nTypeScript", ACCENT),
    ("Backend", "FastAPI\nPython 3.12", GREEN),
    ("DB / Auth / Realtime", "Supabase\n(Postgres · Auth · Storage)", ORANGE),
    ("배포", "Vercel (백엔드)\nExpo EAS (모바일)", NAVY),
]
x0, y0 = Inches(0.8), Inches(1.6)
bw, bh, gap = Inches(2.95), Inches(2.0), Inches(0.15)
for i, (t, d, c) in enumerate(stacks):
    x = x0 + (bw + gap) * i
    add_rect(s, x, y0, bw, Inches(0.6), c)
    add_text(s, x, y0 + Inches(0.13), bw, Inches(0.4),
             t, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, y0 + Inches(0.6), bw, bh - Inches(0.6), LIGHT)
    add_text(s, x + Inches(0.2), y0 + Inches(0.85), bw - Inches(0.4), bh - Inches(0.8),
             d, size=13, color=DARK, align=PP_ALIGN.CENTER)

# Flow diagram (bottom)
fy = Inches(4.2)
add_text(s, Inches(0.8), fy, Inches(12), Inches(0.4),
         "트래픽 분할",
         size=16, bold=True, color=NAVY)

box1 = (Inches(0.8), fy + Inches(0.5), Inches(3.5), Inches(2.0))
box2 = (Inches(4.8), fy + Inches(0.5), Inches(3.5), Inches(2.0))
box3 = (Inches(8.8), fy + Inches(0.5), Inches(3.7), Inches(2.0))
for x, y, w, h in (box1, box2, box3):
    add_rect(s, x, y, w, h, LIGHT)

add_text(s, box1[0], box1[1] + Inches(0.15), box1[2], Inches(0.4),
         "Mobile", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_bullets(s, box1[0] + Inches(0.3), box1[1] + Inches(0.6), box1[2] - Inches(0.4), Inches(1.5),
            ["UI / 화면", "위치 백그라운드 트래커", "Realtime 구독"], size=12)

add_text(s, box2[0], box2[1] + Inches(0.15), box2[2], Inches(0.4),
         "FastAPI", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_bullets(s, box2[0] + Inches(0.3), box2[1] + Inches(0.6), box2[2] - Inches(0.4), Inches(1.5),
            ["권한 검증 · 비즈 로직", "초대 코드 발급/소비", "메시지 전송"], size=12)

add_text(s, box3[0], box3[1] + Inches(0.15), box3[2], Inches(0.4),
         "Supabase", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_bullets(s, box3[0] + Inches(0.3), box3[1] + Inches(0.6), box3[2] - Inches(0.4), Inches(1.5),
            ["Postgres + RLS", "Realtime (채팅 · 위치)", "Storage (이미지)"], size=12)

footer(s, 7)

# ── Slide 8: 프라이버시·보안 ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "프라이버시 & 보안", "친구 사이 신뢰가 곧 제품 가치")

items = [
    ("위치 공유는 약속 시간에만",
     "약속 시작 N분 전(기본 20분) 자동 ON, 종료 시각에 자동 OFF. 상시 추적 없음."),
    ("위치 데이터 단명",
     "위치 핑은 약속 종료 + 24시간 후 DB에서 자동 삭제 (pg_cron)."),
    ("초대 기반 친구 관계",
     "전체 사용자 검색 불가. 초대 링크 또는 QR로만 친구 추가 — 모르는 사람 노출 차단."),
    ("Row Level Security",
     "Supabase RLS 전면 적용. 그룹·약속·채팅·위치 모두 멤버/참여자만 접근 가능."),
    ("소셜 로그인",
     "비밀번호 보관 없음 (Google · Apple · Kakao OAuth). 개인정보 최소 수집."),
]
y = Inches(1.5)
for t, d in items:
    add_rect(s, Inches(0.8), y, Inches(0.15), Inches(1.0), GREEN)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(11), Inches(0.5),
             t, size=16, bold=True, color=DARK)
    add_text(s, Inches(1.2), y + Inches(0.5), Inches(11), Inches(0.5),
             d, size=13, color=MUTED)
    y += Inches(1.1)
footer(s, 8)

# ── Slide 9: 일정 / 마일스톤 ────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "마일스톤 (예상)", "MVP 우선, 단계적 확장")

phases = [
    ("Phase 0", "설계 확정", "스펙 · 데이터 모델 · RLS 정책 확정", ACCENT),
    ("Phase 1", "인증 + 그룹/친구", "소셜 로그인 · 초대 · 그룹 관리", GREEN),
    ("Phase 2", "약속 + 알림", "약속 CRUD · 개인 알림 · 푸시", ORANGE),
    ("Phase 3", "채팅", "그룹/약속 채팅 · 이미지 · 장소 카드", NAVY),
    ("Phase 4", "위치 공유", "백그라운드 트래커 · 지도 · Realtime", ACCENT),
    ("Phase 5", "베타 테스트", "내부 dogfooding → 친구 그룹 클로즈드 베타", GREEN),
]
x0, y0 = Inches(0.8), Inches(1.7)
bw, gap = Inches(2.0), Inches(0.05)
for i, (p, t, d, c) in enumerate(phases):
    x = x0 + (bw + gap) * i
    add_rect(s, x, y0, bw, Inches(0.55), c)
    add_text(s, x, y0 + Inches(0.13), bw, Inches(0.35),
             p, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, y0 + Inches(0.55), bw, Inches(3.5), LIGHT)
    add_text(s, x + Inches(0.15), y0 + Inches(0.7), bw - Inches(0.3), Inches(0.6),
             t, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y0 + Inches(1.5), bw - Inches(0.3), Inches(2.5),
             d, size=11, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.7), LIGHT)
add_text(s, Inches(1.0), Inches(6.18), Inches(11.5), Inches(0.4),
         "각 Phase는 독립 배포 가능 단위. 일정은 설계 검토 후 implementation plan에서 확정.",
         size=13, color=DARK)
footer(s, 9)

# ── Slide 10: 다음 단계 / 논의 ───────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, WHITE)
header(s, "다음 단계 & 논의 사항", "함께 정해야 할 것들")

# Left: next steps
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8), LIGHT)
add_rect(s, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.6), GREEN)
add_text(s, Inches(0.8), Inches(1.73), Inches(5.8), Inches(0.4),
         "✅  다음 단계", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(1.1), Inches(2.4), Inches(5.2), Inches(3.8), [
    "스펙 문서 리뷰 및 피드백 수렴",
    "Implementation plan 작성 (Phase별 태스크)",
    "Supabase 신규 프로젝트 생성",
    "Phase 1 (인증 + 그룹) 착수",
    "내부 dogfooding 그룹 모집",
], size=14)

# Right: open questions
add_rect(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8), LIGHT)
add_rect(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.6), ORANGE)
add_text(s, Inches(6.8), Inches(1.73), Inches(5.8), Inches(0.4),
         "💬  논의가 필요해요", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(3.8), [
    "브랜드 톤 / 로고 방향",
    "베타 테스트 대상 그룹",
    "푸시 알림 빈도/문구 가이드",
    "iOS '항상 허용' 거부 시 UX",
    "출시 후 운영 체계 (이슈 트래킹)",
], size=14)

# Closing
add_text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
         "스펙 문서: MeetPod/docs/superpowers/specs/2026-05-02-meetpod-mvp-design.md",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
footer(s, 10)

out = "d:/Workspace/CPKWorks/MeetPod/docs/MeetPod_프로젝트_소개.pptx"
prs.save(out)
print(f"saved: {out}")
