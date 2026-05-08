"""MeetPod 화면 설계 PPT — iPhone mockup."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Design tokens
NAVY = RGBColor(0x0F, 0x1B, 0x3D)
BRAND = RGBColor(0x5B, 0x7C, 0xFA)        # primary blue/purple
BRAND_2 = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)        # cyan accent
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARN = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_2 = RGBColor(0x33, 0x40, 0x5C)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
HAIR = RGBColor(0xE2, 0xE8, 0xF0)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HERMES = RGBColor(0xFF, 0x69, 0x00)        # 에르메스 오렌지
APP_BG = HERMES                            # 앱 화면 기본 배경
BLACK = RGBColor(0x00, 0x00, 0x00)
PILL_BG = RGBColor(0xEE, 0xF2, 0xFF)
GRAY_BG = RGBColor(0xF1, 0xF5, 0xF9)
BUBBLE_OTHER = RGBColor(0xE5, 0xE7, 0xEB)
BUBBLE_ME = RGBColor(0x5B, 0x7C, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "맑은 고딕"


# ─── helpers ────────────────────────────────────────────────
def shp_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def shp_round(slide, x, y, w, h, fill, line=None, radius_ratio=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if radius_ratio is not None:
        s.adjustments[0] = radius_ratio
    return s


def shp_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, t, size=11, bold=False, color=INK,
        align=PP_ALIGN.LEFT, vcenter=False, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if vcenter:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def line(slide, x1, y1, x2, y2, color=HAIR, weight=0.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    return c


# ─── iPhone frame ──────────────────────────────────────────
# iPhone 14/15 프레임. 화면 가용 영역은 (px, py, pw, ph) 반환.
def iphone(slide, x, y, scale=1.0):
    """iPhone 외곽선 + status bar + dynamic island + home indicator.
    리턴: (inner_x, inner_y, inner_w, inner_h) — 콘텐츠 그릴 영역."""
    W = Inches(2.85 * scale)
    H = Inches(5.85 * scale)
    R = 0.10  # corner radius ratio
    # outer body
    body = shp_round(slide, x, y, W, H, BLACK, radius_ratio=R)
    # bezel inner
    bezel_pad = Inches(0.06 * scale)
    inner_x = x + bezel_pad
    inner_y = y + bezel_pad
    inner_w = W - bezel_pad * 2
    inner_h = H - bezel_pad * 2
    shp_round(slide, inner_x, inner_y, inner_w, inner_h, WHITE,
              radius_ratio=R - 0.02)
    # dynamic island
    island_w = Inches(0.85 * scale)
    island_h = Inches(0.20 * scale)
    island_x = x + (W - island_w) / 2
    island_y = inner_y + Inches(0.10 * scale)
    shp_round(slide, island_x, island_y, island_w, island_h, BLACK,
              radius_ratio=0.5)
    # status bar text
    sx = inner_x + Inches(0.18 * scale)
    sy = inner_y + Inches(0.13 * scale)
    txt(slide, sx, sy, Inches(0.6 * scale), Inches(0.18 * scale),
        "9:41", size=8, bold=True, color=INK)
    # right status icons (signal/wifi/battery as rectangles)
    rx = x + W - Inches(0.5 * scale)
    ry = sy + Inches(0.04 * scale)
    shp_rect(slide, rx, ry, Inches(0.08 * scale), Inches(0.07 * scale), INK)
    shp_rect(slide, rx + Inches(0.10 * scale), ry,
             Inches(0.08 * scale), Inches(0.07 * scale), INK)
    shp_rect(slide, rx + Inches(0.20 * scale), ry,
             Inches(0.18 * scale), Inches(0.08 * scale), INK)
    # content area: below status, above home indicator
    content_top = inner_y + Inches(0.42 * scale)
    content_bottom = inner_y + inner_h - Inches(0.20 * scale)
    # home indicator
    hx = x + W / 2 - Inches(0.55 * scale)
    hy = inner_y + inner_h - Inches(0.12 * scale)
    shp_round(slide, hx, hy, Inches(1.10 * scale), Inches(0.05 * scale),
              INK, radius_ratio=0.5)
    return (inner_x, content_top, inner_w,
            content_bottom - content_top, x, y, W, H)


def avatar(slide, x, y, d, color, initial=""):
    shp_oval(slide, x, y, d, d, color)
    if initial:
        txt(slide, x, y + d * 0.18, d, d * 0.6,
            initial, size=int(d.inches * 22),
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def pill(slide, x, y, w, h, label, fill=PILL_BG, color=BRAND, size=8):
    shp_round(slide, x, y, w, h, fill, radius_ratio=0.5)
    txt(slide, x, y + h * 0.18, w, h * 0.7,
        label, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def fab(slide, cx, cy, d, color=BRAND, label="+"):
    x = cx - d / 2; y = cy - d / 2
    # shadow
    shp_oval(slide, x + Inches(0.02), y + Inches(0.04), d, d,
             RGBColor(0xCB, 0xD5, 0xE1))
    shp_oval(slide, x, y, d, d, color)
    txt(slide, x, y + d * 0.20, d, d * 0.6,
        label, size=int(d.inches * 26), bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)


def tabbar(slide, ix, iy, iw, ih, active=0):
    """하단 탭바: Meetups / Groups / Chats / Me"""
    bar_h = Inches(0.40)
    by = iy + ih - bar_h
    shp_rect(slide, ix, by, iw, Inches(0.005), HAIR)
    shp_rect(slide, ix, by + Inches(0.005), iw, bar_h - Inches(0.005), WHITE)
    items = ["📅", "👥", "💬", "👤"]
    labels = ["약속", "그룹", "채팅", "Me"]
    cw = iw / 4
    for i in range(4):
        col = BRAND if i == active else MUTED
        cx = ix + cw * i
        txt(slide, cx, by + Inches(0.05), cw, Inches(0.20),
            items[i], size=12, color=col, align=PP_ALIGN.CENTER)
        txt(slide, cx, by + Inches(0.24), cw, Inches(0.14),
            labels[i], size=7, bold=(i == active),
            color=col, align=PP_ALIGN.CENTER)


def navbar(slide, ix, iy, iw, title, back=False, action=None):
    """상단 네비: 뒤로가기 + 타이틀 + 우측 액션"""
    h = Inches(0.45)
    shp_rect(slide, ix, iy, iw, h, WHITE)
    shp_rect(slide, ix, iy + h, iw, Inches(0.005), HAIR)
    if back:
        txt(slide, ix + Inches(0.1), iy + Inches(0.10),
            Inches(0.4), Inches(0.3),
            "‹", size=18, bold=True, color=BRAND)
    txt(slide, ix, iy + Inches(0.13), iw, Inches(0.3),
        title, size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)
    if action:
        txt(slide, ix + iw - Inches(0.6), iy + Inches(0.14),
            Inches(0.5), Inches(0.3),
            action, size=10, bold=True, color=BRAND, align=PP_ALIGN.RIGHT)
    return iy + h


# ─── slide chrome ───────────────────────────────────────────
def slide_header(slide, num, title, desc):
    shp_rect(slide, 0, 0, SW, Inches(0.85), NAVY)
    shp_rect(slide, 0, Inches(0.85), SW, Inches(0.04), BRAND)
    txt(slide, Inches(0.4), Inches(0.13),
        Inches(0.7), Inches(0.5),
        f"{num:02d}", size=22, bold=True, color=BRAND)
    txt(slide, Inches(1.1), Inches(0.13),
        Inches(11), Inches(0.4),
        title, size=20, bold=True, color=WHITE)
    txt(slide, Inches(1.1), Inches(0.50),
        Inches(11), Inches(0.3),
        desc, size=11, color=HAIR)


def slide_footer(slide, page, total):
    txt(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
        "MeetPod — 화면 설계", size=9, color=MUTED)
    txt(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
        f"{page} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ─── annotation panel (right side of iPhone) ───────────────
def panel(slide, x, y, w, h, title, color=BRAND):
    shp_round(slide, x, y, w, h, SURFACE, radius_ratio=0.05)
    shp_rect(slide, x, y, Inches(0.08), h, color)
    txt(slide, x + Inches(0.25), y + Inches(0.12),
        w - Inches(0.4), Inches(0.3),
        title, size=12, bold=True, color=INK)


def panel_items(slide, x, y, w, items, size=10):
    for i, it in enumerate(items):
        cy = y + Inches(0.05 + i * 0.32)
        shp_oval(slide, x, cy + Inches(0.07),
                 Inches(0.08), Inches(0.08), BRAND)
        txt(slide, x + Inches(0.18), cy,
            w - Inches(0.2), Inches(0.3),
            it, size=size, color=INK_2)


# ============================================================
# Page count
# ============================================================
TOTAL = 16  # cover + design system + nav + 12 screens + summary

# ─── 1. Cover ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
shp_rect(s, 0, 0, SW, SH, NAVY)
# diagonal accent stripe
shp_rect(s, 0, Inches(3.2), SW, Inches(0.05), BRAND)
shp_rect(s, 0, Inches(3.30), SW, Inches(0.02), BRAND_2)
txt(s, Inches(0.8), Inches(1.6), Inches(12), Inches(1.2),
    "MeetPod", size=64, bold=True, color=WHITE)
txt(s, Inches(0.8), Inches(2.6), Inches(12), Inches(0.5),
    "화면 설계 (UI Design)", size=28, bold=True, color=BRAND)
txt(s, Inches(0.8), Inches(3.4), Inches(12), Inches(0.5),
    "iPhone 프레임 기반 12개 핵심 화면 + 디자인 시스템",
    size=16, color=HAIR)

# little phone illustration
mini_x = Inches(10.0); mini_y = Inches(4.5)
shp_round(s, mini_x, mini_y, Inches(1.8), Inches(2.6),
          BLACK, radius_ratio=0.10)
shp_round(s, mini_x + Inches(0.07), mini_y + Inches(0.07),
          Inches(1.66), Inches(2.46), BRAND, radius_ratio=0.08)
shp_round(s, mini_x + Inches(0.65), mini_y + Inches(0.18),
          Inches(0.5), Inches(0.12), BLACK, radius_ratio=0.5)
txt(s, mini_x, mini_y + Inches(1.2), Inches(1.8), Inches(0.4),
    "MeetPod", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
    "2026.05.03  ·  CPKWorks", size=11, color=MUTED)

# ─── 2. Design System ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
shp_rect(s, 0, 0, SW, SH, WHITE)
slide_header(s, 1, "Design System",
             "컬러 팔레트 · 타이포그래피 · 컴포넌트 기본")

# Colors
txt(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
    "Colors", size=14, bold=True, color=INK)
swatches = [
    ("App BG (Hermès)", HERMES, "#FF6900"),
    ("Brand", BRAND, "#5B7CFA"),
    ("Brand 2", BRAND_2, "#8B5CF6"),
    ("Accent", ACCENT, "#06B6D4"),
    ("Success", SUCCESS, "#10B981"),
    ("Warn", WARN, "#F59E0B"),
    ("Danger", DANGER, "#EF4444"),
    ("Ink", INK, "#0F172A"),
]
for i, (name, c, hex_) in enumerate(swatches):
    col, row = i % 4, i // 4
    x = Inches(0.5 + col * 1.55)
    y = Inches(1.55 + row * 1.10)
    shp_round(s, x, y, Inches(1.4), Inches(0.7), c, radius_ratio=0.10)
    txt(s, x, y + Inches(0.78), Inches(1.4), Inches(0.25),
        name, size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, x, y + Inches(0.99), Inches(1.4), Inches(0.2),
        hex_, size=8, color=MUTED, align=PP_ALIGN.CENTER)

# Typography
txt(s, Inches(7.0), Inches(1.1), Inches(6), Inches(0.4),
    "Typography", size=14, bold=True, color=INK)
typo = [
    ("Display", 26, True, "약속을 더 쉽게"),
    ("Title", 18, True, "오늘의 약속"),
    ("Heading", 14, True, "참여자 (3)"),
    ("Body", 11, False, "친구 두 명과 7시 강남역에서 만나기"),
    ("Caption", 9, False, "20분 전 자동 시작"),
]
for i, (name, sz, b, sample) in enumerate(typo):
    y = Inches(1.55 + i * 0.55)
    txt(s, Inches(7.0), y, Inches(1.5), Inches(0.4),
        f"{name} · {sz}pt", size=9, color=MUTED)
    txt(s, Inches(8.5), y - Inches(0.05), Inches(4.5), Inches(0.5),
        sample, size=sz, bold=b, color=INK)

# Components
txt(s, Inches(0.5), Inches(4.4), Inches(6), Inches(0.4),
    "Components", size=14, bold=True, color=INK)
# button primary
shp_round(s, Inches(0.5), Inches(4.85), Inches(2.0), Inches(0.45),
          BRAND, radius_ratio=0.5)
txt(s, Inches(0.5), Inches(4.95), Inches(2.0), Inches(0.3),
    "Primary 버튼", size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# button secondary
shp_round(s, Inches(2.7), Inches(4.85), Inches(2.0), Inches(0.45),
          WHITE, line=BRAND, radius_ratio=0.5)
txt(s, Inches(2.7), Inches(4.95), Inches(2.0), Inches(0.3),
    "Secondary 버튼", size=10, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
# input
shp_round(s, Inches(0.5), Inches(5.45), Inches(4.2), Inches(0.5),
          GRAY_BG, radius_ratio=0.25)
txt(s, Inches(0.7), Inches(5.6), Inches(3.8), Inches(0.3),
    "닉네임을 입력하세요", size=10, color=MUTED)
# pill
pill(s, Inches(0.5), Inches(6.1), Inches(1.0), Inches(0.3),
     "그룹", PILL_BG, BRAND, 9)
pill(s, Inches(1.55), Inches(6.1), Inches(1.0), Inches(0.3),
     "약속", RGBColor(0xDC, 0xFC, 0xE7), SUCCESS, 9)
pill(s, Inches(2.6), Inches(6.1), Inches(1.0), Inches(0.3),
     "긴급", RGBColor(0xFE, 0xE2, 0xE2), DANGER, 9)

# message bubbles
shp_round(s, Inches(7.0), Inches(4.85), Inches(2.4), Inches(0.5),
          BUBBLE_OTHER, radius_ratio=0.4)
txt(s, Inches(7.15), Inches(4.97), Inches(2.2), Inches(0.3),
    "어디야? 곧 도착해", size=10, color=INK)
shp_round(s, Inches(8.6), Inches(5.5), Inches(2.4), Inches(0.5),
          BUBBLE_ME, radius_ratio=0.4)
txt(s, Inches(8.75), Inches(5.62), Inches(2.2), Inches(0.3),
    "5분 안에 도착!", size=10, color=WHITE)
# place card mini
shp_round(s, Inches(7.0), Inches(6.15), Inches(4.0), Inches(0.6),
          WHITE, line=HAIR, radius_ratio=0.15)
shp_round(s, Inches(7.1), Inches(6.22), Inches(0.45), Inches(0.45),
          PILL_BG, radius_ratio=0.3)
txt(s, Inches(7.1), Inches(6.27), Inches(0.45), Inches(0.4),
    "📍", size=14, align=PP_ALIGN.CENTER)
txt(s, Inches(7.65), Inches(6.20), Inches(3.3), Inches(0.25),
    "OO식당 강남점", size=9, bold=True, color=INK)
txt(s, Inches(7.65), Inches(6.42), Inches(3.3), Inches(0.25),
    "서울 강남구 테헤란로 123", size=8, color=MUTED)

slide_footer(s, 2, TOTAL)

# ─── 3. Navigation Map ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
shp_rect(s, 0, 0, SW, SH, WHITE)
slide_header(s, 2, "Navigation Map", "앱 전체 화면 구조")

# Auth flow (top)
shp_round(s, Inches(0.5), Inches(1.2), Inches(4), Inches(0.7),
          NAVY, radius_ratio=0.2)
txt(s, Inches(0.5), Inches(1.35), Inches(4), Inches(0.4),
    "Auth Stack", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

auth_screens = ["LoginScreen", "OnboardingHandleScreen", "InviteAcceptScreen"]
for i, name in enumerate(auth_screens):
    x = Inches(0.5 + i * 1.4)
    shp_round(s, x, Inches(2.05), Inches(1.30), Inches(0.55),
              SURFACE, line=BRAND, radius_ratio=0.2)
    txt(s, x, Inches(2.18), Inches(1.30), Inches(0.4),
        name.replace("Screen", ""), size=8, bold=True,
        color=BRAND, align=PP_ALIGN.CENTER, vcenter=True)

# Main tab (bottom)
shp_round(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(4.0),
          SURFACE, line=HAIR, radius_ratio=0.05)
shp_round(s, Inches(0.5), Inches(3.0), Inches(4.0), Inches(0.7),
          NAVY, radius_ratio=0.2)
txt(s, Inches(0.5), Inches(3.15), Inches(4.0), Inches(0.4),
    "Main Tab Navigator", size=14, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

tabs = [
    ("📅 Meetups",
     ["MeetupList", "MeetupDetail", "MeetupCreate",
      "MeetupMap", "PlacePicker"], BRAND),
    ("👥 Groups",
     ["GroupList", "GroupDetail", "GroupCreate",
      "GroupMembers", "GroupInvite"], SUCCESS),
    ("💬 Chats",
     ["ChatList", "ChatRoom"], WARN),
    ("👤 Me",
     ["Profile", "RemindersDefault"], BRAND_2),
]
tx0 = Inches(0.7)
tw = Inches(2.95)
gap = Inches(0.07)
for i, (name, screens, color) in enumerate(tabs):
    x = tx0 + (tw + gap) * i
    shp_round(s, x, Inches(3.95), tw, Inches(0.55),
              color, radius_ratio=0.2)
    txt(s, x, Inches(4.08), tw, Inches(0.4),
        name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, sn in enumerate(screens):
        sy = Inches(4.65 + j * 0.50)
        shp_round(s, x + Inches(0.15), sy, tw - Inches(0.3),
                  Inches(0.42), WHITE, line=color, radius_ratio=0.3)
        txt(s, x + Inches(0.15), sy + Inches(0.10),
            tw - Inches(0.3), Inches(0.3),
            sn, size=9, bold=True, color=INK,
            align=PP_ALIGN.CENTER, vcenter=True)

slide_footer(s, 3, TOTAL)


# ============================================================
# Screen slide template
# ============================================================
def screen_slide(page, num, name_kr, name_en, route, draw_screen,
                 layout_notes, interaction_notes, design_notes):
    """draw_screen(slide, ix, iy, iw, ih) — iPhone 내부에 그릴 함수"""
    s = prs.slides.add_slide(BLANK)
    shp_rect(s, 0, 0, SW, SH, WHITE)
    slide_header(s, num, f"{name_kr}  ·  {name_en}",
                 f"Route: {route}")
    # iPhone (좌측)
    px = Inches(0.7); py = Inches(1.15)
    ix, iy, iw, ih, _, _, _, _ = iphone(s, px, py, scale=1.05)
    draw_screen(s, ix, iy, iw, ih)
    # 우측 패널
    rx = Inches(4.6); rw = Inches(8.2)
    panel(s, rx, Inches(1.15), rw, Inches(1.85), "Layout", BRAND)
    panel_items(s, rx + Inches(0.3), Inches(1.55), rw - Inches(0.5),
                layout_notes)
    panel(s, rx, Inches(3.10), rw, Inches(1.85), "Interaction", SUCCESS)
    panel_items(s, rx + Inches(0.3), Inches(3.50), rw - Inches(0.5),
                interaction_notes)
    panel(s, rx, Inches(5.05), rw, Inches(2.0), "Design Note", BRAND_2)
    panel_items(s, rx + Inches(0.3), Inches(5.45), rw - Inches(0.5),
                design_notes)
    slide_footer(s, page, TOTAL)


# ============================================================
# Draw functions for each screen
# ============================================================

def draw_login(s, ix, iy, iw, ih):
    # Background gradient (use solid + accent stripe)
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    # logo area
    cy = iy + Inches(0.5)
    shp_round(s, ix + iw / 2 - Inches(0.45), cy,
              Inches(0.9), Inches(0.9), BRAND, radius_ratio=0.3)
    txt(s, ix + iw / 2 - Inches(0.45), cy + Inches(0.20),
        Inches(0.9), Inches(0.5),
        "M", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ix, cy + Inches(1.05), iw, Inches(0.4),
        "MeetPod", size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, ix, cy + Inches(1.4), iw, Inches(0.3),
        "친구들의 약속을 더 쉽게", size=10,
        color=MUTED, align=PP_ALIGN.CENTER)
    # buttons
    by = iy + Inches(3.1)
    btn_w = iw - Inches(0.5)
    btn_x = ix + Inches(0.25)
    socials = [
        ("Apple로 계속하기", BLACK, WHITE),
        ("Google로 계속하기", WHITE, INK),
        ("Kakao로 계속하기", RGBColor(0xFE, 0xE5, 0x00), INK),
    ]
    for i, (label, bg, fg) in enumerate(socials):
        y = by + Inches(i * 0.58)
        line_color = HAIR if bg == WHITE else None
        shp_round(s, btn_x, y, btn_w, Inches(0.48), bg, line=line_color,
                  radius_ratio=0.5)
        txt(s, btn_x, y + Inches(0.13), btn_w, Inches(0.3),
            label, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
    # legal
    txt(s, ix, iy + ih - Inches(0.65), iw, Inches(0.3),
        "계속 진행 시 이용약관 · 개인정보 처리방침 동의",
        size=7, color=MUTED, align=PP_ALIGN.CENTER)


def draw_onboarding(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "프로필 설정")
    # avatar upload
    cy = ny + Inches(0.4)
    shp_oval(s, ix + iw / 2 - Inches(0.55), cy,
             Inches(1.10), Inches(1.10), GRAY_BG)
    txt(s, ix + iw / 2 - Inches(0.55), cy + Inches(0.30),
        Inches(1.10), Inches(0.5),
        "📷", size=22, color=MUTED, align=PP_ALIGN.CENTER)
    txt(s, ix, cy + Inches(1.20), iw, Inches(0.3),
        "사진 추가", size=9, color=BRAND, align=PP_ALIGN.CENTER)
    # handle
    fy = cy + Inches(1.7)
    txt(s, ix + Inches(0.3), fy, iw, Inches(0.3),
        "핸들 (변경 불가)", size=9, bold=True, color=INK)
    shp_round(s, ix + Inches(0.25), fy + Inches(0.32),
              iw - Inches(0.5), Inches(0.5), GRAY_BG, radius_ratio=0.3)
    txt(s, ix + Inches(0.4), fy + Inches(0.45),
        Inches(0.3), Inches(0.3),
        "@", size=14, bold=True, color=MUTED)
    txt(s, ix + Inches(0.7), fy + Inches(0.47),
        iw - Inches(1.0), Inches(0.3),
        "hyejin22", size=12, color=INK)
    txt(s, ix + Inches(0.3), fy + Inches(0.9), iw, Inches(0.3),
        "✓ 사용 가능한 핸들이에요", size=8, color=SUCCESS)
    # display name
    fy2 = fy + Inches(1.3)
    txt(s, ix + Inches(0.3), fy2, iw, Inches(0.3),
        "표시 이름", size=9, bold=True, color=INK)
    shp_round(s, ix + Inches(0.25), fy2 + Inches(0.32),
              iw - Inches(0.5), Inches(0.5), GRAY_BG, radius_ratio=0.3)
    txt(s, ix + Inches(0.4), fy2 + Inches(0.47),
        iw - Inches(0.7), Inches(0.3),
        "혜진", size=12, color=INK)
    # CTA
    by = iy + ih - Inches(0.95)
    shp_round(s, ix + Inches(0.25), by, iw - Inches(0.5), Inches(0.5),
              BRAND, radius_ratio=0.5)
    txt(s, ix + Inches(0.25), by + Inches(0.13),
        iw - Inches(0.5), Inches(0.3),
        "시작하기", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def draw_meetup_list(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    # large title nav
    h = Inches(0.85)
    shp_rect(s, ix, iy, iw, h, WHITE)
    txt(s, ix + Inches(0.3), iy + Inches(0.2), iw, Inches(0.5),
        "약속", size=22, bold=True, color=INK)
    # search
    sy = iy + h + Inches(0.1)
    shp_round(s, ix + Inches(0.25), sy, iw - Inches(0.5), Inches(0.42),
              GRAY_BG, radius_ratio=0.5)
    txt(s, ix + Inches(0.4), sy + Inches(0.11),
        iw - Inches(0.7), Inches(0.3),
        "🔍  약속 검색", size=10, color=MUTED)
    # segment
    seg_y = sy + Inches(0.55)
    shp_round(s, ix + Inches(0.25), seg_y, iw - Inches(0.5), Inches(0.36),
              GRAY_BG, radius_ratio=0.5)
    seg_w = (iw - Inches(0.5)) / 3
    shp_round(s, ix + Inches(0.28), seg_y + Inches(0.03),
              seg_w - Inches(0.06), Inches(0.30), WHITE, radius_ratio=0.5)
    for i, lbl in enumerate(["다가오는", "오늘", "지난"]):
        sx = ix + Inches(0.25) + seg_w * i
        col = INK if i == 0 else MUTED
        b = (i == 0)
        txt(s, sx, seg_y + Inches(0.09), seg_w, Inches(0.25),
            lbl, size=9, bold=b, color=col, align=PP_ALIGN.CENTER)
    # cards
    cy = seg_y + Inches(0.55)
    cards = [
        ("오늘 · 19:00", "강남역 저녁",
         "OO식당 · 강남역 3번출구", "지수 · +2", BRAND, "📍 20분 후 위치공유"),
        ("내일 · 12:30", "북카페 모임",
         "동네책방 · 망원동", "민호 · +3", SUCCESS, "🔔 30분 전 알림"),
        ("토 · 06:00", "OO산악회 산행",
         "수락산 입구", "태훈 · +11", BRAND_2, "📅 그룹 약속"),
    ]
    for i, (when, title, place, host, color, badge) in enumerate(cards):
        y = cy + Inches(i * 1.15)
        shp_round(s, ix + Inches(0.25), y, iw - Inches(0.5),
                  Inches(1.05), WHITE, line=HAIR, radius_ratio=0.10)
        # left color stripe
        shp_round(s, ix + Inches(0.25), y, Inches(0.10),
                  Inches(1.05), color, radius_ratio=0.3)
        txt(s, ix + Inches(0.45), y + Inches(0.08),
            iw - Inches(0.7), Inches(0.25),
            when, size=8, bold=True, color=color)
        txt(s, ix + Inches(0.45), y + Inches(0.30),
            iw - Inches(0.7), Inches(0.3),
            title, size=12, bold=True, color=INK)
        txt(s, ix + Inches(0.45), y + Inches(0.55),
            iw - Inches(0.7), Inches(0.25),
            place, size=8, color=MUTED)
        txt(s, ix + Inches(0.45), y + Inches(0.75),
            iw - Inches(0.7), Inches(0.25),
            host, size=8, color=INK_2)
        # badge
        pill(s, ix + iw - Inches(1.7), y + Inches(0.75),
             Inches(1.4), Inches(0.22), badge, PILL_BG, color, 7)
    # FAB
    fab(s, ix + iw - Inches(0.6), iy + ih - Inches(0.95),
        Inches(0.65), BRAND, "+")
    tabbar(s, ix, iy, iw, ih, active=0)


def draw_meetup_detail(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "약속 상세", back=True, action="•••")
    # hero
    hy = ny
    shp_rect(s, ix, hy, iw, Inches(1.4), BRAND)
    txt(s, ix + Inches(0.3), hy + Inches(0.2), iw, Inches(0.3),
        "오늘 19:00 · 친구 약속", size=9, color=HAIR)
    txt(s, ix + Inches(0.3), hy + Inches(0.45), iw - Inches(0.6),
        Inches(0.5),
        "강남역 저녁", size=20, bold=True, color=WHITE)
    txt(s, ix + Inches(0.3), hy + Inches(0.95), iw - Inches(0.6),
        Inches(0.4),
        "📍 OO식당 · 강남역 3번출구 도보 2분",
        size=9, color=HAIR)
    # info rows
    iy0 = hy + Inches(1.55)
    rows = [
        ("👥", "참여자", "지수 · 민호 · 서연 (3)"),
        ("⏰", "내 알림", "30분 전 · 변경"),
        ("📡", "위치 공유", "18:40 자동 시작 (20분 전)"),
    ]
    for i, (icon, k, v) in enumerate(rows):
        y = iy0 + Inches(i * 0.55)
        shp_round(s, ix + Inches(0.25), y, iw - Inches(0.5),
                  Inches(0.48), WHITE, line=HAIR, radius_ratio=0.15)
        txt(s, ix + Inches(0.35), y + Inches(0.13),
            Inches(0.3), Inches(0.3),
            icon, size=12, color=BRAND)
        txt(s, ix + Inches(0.7), y + Inches(0.07),
            Inches(1.0), Inches(0.2),
            k, size=8, color=MUTED)
        txt(s, ix + Inches(0.7), y + Inches(0.24),
            iw - Inches(1.0), Inches(0.25),
            v, size=10, bold=True, color=INK)
    # avatars
    ay = iy0 + Inches(1.85)
    txt(s, ix + Inches(0.3), ay, iw, Inches(0.3),
        "지도에서 위치 보기", size=10, bold=True, color=INK)
    shp_round(s, ix + Inches(0.25), ay + Inches(0.32),
              iw - Inches(0.5), Inches(1.0),
              RGBColor(0xDB, 0xEA, 0xFE), radius_ratio=0.10)
    # mini map dots
    for i, (mx, my, c) in enumerate([
        (0.6, 0.5, BRAND), (1.4, 0.7, SUCCESS), (2.0, 0.4, WARN)
    ]):
        shp_oval(s, ix + Inches(mx), ay + Inches(0.32 + my * 0.3),
                 Inches(0.18), Inches(0.18), c, line=WHITE)
    # bottom CTA
    by = iy + ih - Inches(1.0)
    shp_round(s, ix + Inches(0.25), by, iw - Inches(0.5), Inches(0.5),
              BRAND, radius_ratio=0.5)
    txt(s, ix + Inches(0.25), by + Inches(0.13),
        iw - Inches(0.5), Inches(0.3),
        "💬 약속 채팅 열기", size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)


def draw_meetup_create(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "새 약속", back=True, action="저장")
    y = ny + Inches(0.15)
    # title input
    shp_round(s, ix + Inches(0.25), y, iw - Inches(0.5), Inches(0.55),
              GRAY_BG, radius_ratio=0.25)
    txt(s, ix + Inches(0.4), y + Inches(0.16), iw, Inches(0.3),
        "약속 제목", size=11, color=MUTED)
    # group selector
    y2 = y + Inches(0.7)
    shp_round(s, ix + Inches(0.25), y2, iw - Inches(0.5), Inches(0.5),
              SURFACE, line=HAIR, radius_ratio=0.2)
    txt(s, ix + Inches(0.4), y2 + Inches(0.13), Inches(2.0), Inches(0.3),
        "그룹", size=9, color=MUTED)
    txt(s, ix + iw - Inches(1.6), y2 + Inches(0.13),
        Inches(1.3), Inches(0.3),
        "1회성  ›", size=10, bold=True, color=BRAND, align=PP_ALIGN.RIGHT)
    # date/time
    y3 = y2 + Inches(0.65)
    for i, (k, v) in enumerate([("시작", "5/15 토 18:00"), ("종료", "5/15 토 21:00")]):
        yy = y3 + Inches(i * 0.5)
        shp_round(s, ix + Inches(0.25), yy, iw - Inches(0.5), Inches(0.42),
                  SURFACE, line=HAIR, radius_ratio=0.2)
        txt(s, ix + Inches(0.4), yy + Inches(0.10), Inches(2.0),
            Inches(0.3), k, size=9, color=MUTED)
        txt(s, ix + iw - Inches(2.5), yy + Inches(0.10),
            Inches(2.2), Inches(0.3),
            v, size=10, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    # place
    y4 = y3 + Inches(1.05)
    shp_round(s, ix + Inches(0.25), y4, iw - Inches(0.5), Inches(0.7),
              SURFACE, line=HAIR, radius_ratio=0.2)
    txt(s, ix + Inches(0.4), y4 + Inches(0.10),
        Inches(2.0), Inches(0.3),
        "📍 장소", size=9, color=MUTED)
    txt(s, ix + Inches(0.4), y4 + Inches(0.32),
        iw - Inches(0.8), Inches(0.3),
        "OO식당 강남점", size=11, bold=True, color=INK)
    # members
    y5 = y4 + Inches(0.85)
    shp_round(s, ix + Inches(0.25), y5, iw - Inches(0.5), Inches(0.6),
              SURFACE, line=HAIR, radius_ratio=0.2)
    txt(s, ix + Inches(0.4), y5 + Inches(0.10), Inches(3.0),
        Inches(0.3), "참여자  ·  3", size=9, color=MUTED)
    # mini avatars
    for i, c in enumerate([BRAND, SUCCESS, WARN]):
        avatar(s, ix + Inches(0.4 + i * 0.30), y5 + Inches(0.30),
               Inches(0.28), c, "")
    txt(s, ix + iw - Inches(1.0), y5 + Inches(0.18),
        Inches(0.7), Inches(0.3),
        "편집  ›", size=9, bold=True, color=BRAND, align=PP_ALIGN.RIGHT)
    # location share
    y6 = y5 + Inches(0.75)
    shp_round(s, ix + Inches(0.25), y6, iw - Inches(0.5), Inches(0.42),
              SURFACE, line=HAIR, radius_ratio=0.2)
    txt(s, ix + Inches(0.4), y6 + Inches(0.10), Inches(2.5),
        Inches(0.3), "📡 위치 공유 시작", size=9, color=MUTED)
    txt(s, ix + iw - Inches(2.0), y6 + Inches(0.10),
        Inches(1.7), Inches(0.3),
        "20분 전  ›", size=10, bold=True,
        color=INK, align=PP_ALIGN.RIGHT)
    # personal reminder
    y7 = y6 + Inches(0.55)
    shp_round(s, ix + Inches(0.25), y7, iw - Inches(0.5), Inches(0.42),
              SURFACE, line=HAIR, radius_ratio=0.2)
    txt(s, ix + Inches(0.4), y7 + Inches(0.10), Inches(2.5),
        Inches(0.3), "🔔 내 알림", size=9, color=MUTED)
    txt(s, ix + iw - Inches(2.0), y7 + Inches(0.10),
        Inches(1.7), Inches(0.3),
        "30분 전  ›", size=10, bold=True,
        color=INK, align=PP_ALIGN.RIGHT)


def draw_meetup_map(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, RGBColor(0xE3, 0xEE, 0xFA))
    ny = navbar(s, ix, iy, iw, "지도", back=True)
    # map area
    map_y = ny
    map_h = ih - (ny - iy) - Inches(1.7)
    shp_rect(s, ix, map_y, iw, map_h, RGBColor(0xDB, 0xEA, 0xFE))
    # streets
    for x_off in [0.5, 1.3, 2.0]:
        shp_rect(s, ix + Inches(x_off), map_y, Inches(0.06),
                 map_h, WHITE)
    for y_off in [0.4, 1.2, 2.0, 2.8]:
        shp_rect(s, ix, map_y + Inches(y_off), iw, Inches(0.06), WHITE)
    # destination pin (large, brand)
    pin_x = ix + iw / 2 - Inches(0.18)
    pin_y = map_y + map_h / 2 - Inches(0.4)
    shp_oval(s, pin_x, pin_y, Inches(0.36), Inches(0.36),
             DANGER, line=WHITE)
    txt(s, pin_x, pin_y + Inches(0.06), Inches(0.36), Inches(0.3),
        "📍", size=10, color=WHITE, align=PP_ALIGN.CENTER)
    # member pins
    for x_off, y_off, c, ini in [
        (0.6, 0.6, BRAND, "지"),
        (2.1, 0.5, SUCCESS, "민"),
        (1.0, 2.4, WARN, "서"),
    ]:
        avatar(s, ix + Inches(x_off), map_y + Inches(y_off),
               Inches(0.40), c, ini)
    # bottom sheet
    sy = iy + ih - Inches(1.7)
    shp_round(s, ix, sy, iw, Inches(1.7), WHITE, radius_ratio=0.15)
    # handle
    shp_round(s, ix + iw / 2 - Inches(0.25), sy + Inches(0.08),
              Inches(0.50), Inches(0.06),
              RGBColor(0xCB, 0xD5, 0xE1), radius_ratio=0.5)
    txt(s, ix + Inches(0.3), sy + Inches(0.20), iw, Inches(0.3),
        "참여자 위치", size=11, bold=True, color=INK)
    # rows
    for i, (ini, name, eta, c) in enumerate([
        ("지", "지수", "도착", SUCCESS),
        ("민", "민호", "5분 거리", WARN),
        ("서", "서연", "12분 거리", DANGER),
    ]):
        ry = sy + Inches(0.55 + i * 0.34)
        avatar(s, ix + Inches(0.3), ry, Inches(0.25), c, ini)
        txt(s, ix + Inches(0.7), ry + Inches(0.04),
            Inches(1.5), Inches(0.25),
            name, size=10, bold=True, color=INK)
        txt(s, ix + iw - Inches(1.4), ry + Inches(0.04),
            Inches(1.0), Inches(0.25),
            eta, size=9, color=c, align=PP_ALIGN.RIGHT, bold=True)


def draw_group_list(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    h = Inches(0.85)
    shp_rect(s, ix, iy, iw, h, WHITE)
    txt(s, ix + Inches(0.3), iy + Inches(0.2), iw, Inches(0.5),
        "그룹", size=22, bold=True, color=INK)
    txt(s, ix + iw - Inches(0.7), iy + Inches(0.32),
        Inches(0.4), Inches(0.4),
        "+", size=22, bold=True, color=BRAND)
    # cards
    groups = [
        ("대학동기", "12명 · 다가오는 약속 1", BRAND, "대"),
        ("OO산악회", "12명 · 다가오는 약속 1", SUCCESS, "산"),
        ("육아맘", "4명 · 다가오는 약속 0", WARN, "육"),
        ("치킨집사장모임", "8명", BRAND_2, "치"),
    ]
    cy = iy + h + Inches(0.15)
    for i, (name, sub, c, ini) in enumerate(groups):
        y = cy + Inches(i * 0.95)
        shp_round(s, ix + Inches(0.25), y, iw - Inches(0.5), Inches(0.85),
                  WHITE, line=HAIR, radius_ratio=0.12)
        avatar(s, ix + Inches(0.4), y + Inches(0.13),
               Inches(0.6), c, ini)
        txt(s, ix + Inches(1.1), y + Inches(0.18),
            iw - Inches(1.4), Inches(0.3),
            name, size=12, bold=True, color=INK)
        txt(s, ix + Inches(1.1), y + Inches(0.48),
            iw - Inches(1.4), Inches(0.25),
            sub, size=8, color=MUTED)
        txt(s, ix + iw - Inches(0.5), y + Inches(0.30),
            Inches(0.3), Inches(0.3),
            "›", size=14, color=MUTED)
    tabbar(s, ix, iy, iw, ih, active=1)


def draw_group_detail(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "그룹", back=True, action="•••")
    # header card
    hy = ny + Inches(0.1)
    shp_round(s, ix + Inches(0.25), hy, iw - Inches(0.5), Inches(1.5),
              WHITE, line=HAIR, radius_ratio=0.15)
    avatar(s, ix + iw / 2 - Inches(0.4), hy + Inches(0.18),
           Inches(0.8), SUCCESS, "산")
    txt(s, ix, hy + Inches(1.05), iw, Inches(0.3),
        "OO산악회", size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, ix, hy + Inches(1.30), iw, Inches(0.3),
        "12명 · 방장 태훈", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    # actions row
    ay = hy + Inches(1.7)
    actions = [("📅", "약속 만들기"), ("👥", "멤버"), ("🔗", "초대")]
    for i, (icon, lbl) in enumerate(actions):
        x = ix + Inches(0.25) + (iw - Inches(0.5)) / 3 * i
        w = (iw - Inches(0.5)) / 3 - Inches(0.05)
        shp_round(s, x, ay, w, Inches(0.7), WHITE,
                  line=HAIR, radius_ratio=0.15)
        txt(s, x, ay + Inches(0.10), w, Inches(0.3),
            icon, size=14, color=BRAND, align=PP_ALIGN.CENTER)
        txt(s, x, ay + Inches(0.40), w, Inches(0.25),
            lbl, size=8, bold=True, color=INK, align=PP_ALIGN.CENTER)
    # upcoming meetups
    sy = ay + Inches(0.95)
    txt(s, ix + Inches(0.3), sy, iw, Inches(0.3),
        "다가오는 약속", size=10, bold=True, color=INK)
    shp_round(s, ix + Inches(0.25), sy + Inches(0.35),
              iw - Inches(0.5), Inches(0.85),
              WHITE, line=HAIR, radius_ratio=0.15)
    shp_round(s, ix + Inches(0.25), sy + Inches(0.35),
              Inches(0.10), Inches(0.85), SUCCESS, radius_ratio=0.3)
    txt(s, ix + Inches(0.45), sy + Inches(0.43),
        iw - Inches(0.7), Inches(0.25),
        "토 06:00", size=8, bold=True, color=SUCCESS)
    txt(s, ix + Inches(0.45), sy + Inches(0.62),
        iw - Inches(0.7), Inches(0.3),
        "수락산 산행", size=11, bold=True, color=INK)
    txt(s, ix + Inches(0.45), sy + Inches(0.92),
        iw - Inches(0.7), Inches(0.25),
        "수락산 입구 · 11명 참여", size=8, color=MUTED)
    # chat link
    cy2 = sy + Inches(1.45)
    shp_round(s, ix + Inches(0.25), cy2, iw - Inches(0.5), Inches(0.55),
              BRAND, radius_ratio=0.5)
    txt(s, ix + Inches(0.25), cy2 + Inches(0.16),
        iw - Inches(0.5), Inches(0.3),
        "💬 그룹 채팅 열기", size=11, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)


def draw_chat_list(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    h = Inches(0.85)
    shp_rect(s, ix, iy, iw, h, WHITE)
    txt(s, ix + Inches(0.3), iy + Inches(0.2), iw, Inches(0.5),
        "채팅", size=22, bold=True, color=INK)
    chats = [
        ("대학동기", "민호: 다음주 어디서 만날까?", "오후 2:14",
         BRAND, "대", 3, False),
        ("강남역 저녁", "지수: 5분 안에 도착!", "오후 1:50",
         SUCCESS, "강", 0, False),
        ("OO산악회", "태훈: 사진 공유합니다 📷", "오전 11:32",
         BRAND_2, "산", 1, False),
        ("동네책방 모임 (지난)", "서연: 즐거웠어요~", "어제",
         MUTED, "동", 0, True),
    ]
    sy = iy + h + Inches(0.05)
    for i, (name, last, time, c, ini, badge, archived) in enumerate(chats):
        y = sy + Inches(i * 0.85)
        avatar(s, ix + Inches(0.3), y + Inches(0.13),
               Inches(0.6), c, ini)
        txt(s, ix + Inches(1.0), y + Inches(0.10),
            iw - Inches(2.2), Inches(0.3),
            name, size=11, bold=True,
            color=MUTED if archived else INK)
        txt(s, ix + Inches(1.0), y + Inches(0.40),
            iw - Inches(1.6), Inches(0.3),
            last, size=9, color=MUTED)
        txt(s, ix + iw - Inches(0.95), y + Inches(0.10),
            Inches(0.7), Inches(0.25),
            time, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
        if badge > 0:
            shp_oval(s, ix + iw - Inches(0.55), y + Inches(0.42),
                     Inches(0.30), Inches(0.30), BRAND)
            txt(s, ix + iw - Inches(0.55), y + Inches(0.46),
                Inches(0.30), Inches(0.25),
                str(badge), size=8, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        if archived:
            pill(s, ix + iw - Inches(1.1), y + Inches(0.42),
                 Inches(0.7), Inches(0.22),
                 "지난", GRAY_BG, MUTED, 7)
        line(s, ix + Inches(1.0), y + Inches(0.83),
             ix + iw, y + Inches(0.83), HAIR, 0.5)
    tabbar(s, ix, iy, iw, ih, active=2)


def draw_chat_room(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "강남역 저녁  · 3", back=True, action="ⓘ")
    # info banner
    by = ny + Inches(0.05)
    shp_round(s, ix + Inches(0.25), by, iw - Inches(0.5), Inches(0.4),
              PILL_BG, radius_ratio=0.25)
    txt(s, ix + Inches(0.4), by + Inches(0.10), iw - Inches(0.7),
        Inches(0.25),
        "📡 18:40부터 위치 공유가 시작돼요", size=8,
        bold=True, color=BRAND)
    # messages
    msgs = [
        ("other", "민", BRAND, "거의 다 왔어!"),
        ("other", "서", WARN, "나 5분 늦을 듯..."),
        ("place", None, None, None),
        ("me", None, None, "나도 5분! 식당 안에서 보자"),
    ]
    my = by + Inches(0.55)
    for kind, ini, c, body in msgs:
        if kind == "other":
            avatar(s, ix + Inches(0.25), my + Inches(0.05),
                   Inches(0.30), c, ini)
            shp_round(s, ix + Inches(0.65), my,
                      Inches(1.7), Inches(0.42),
                      BUBBLE_OTHER, radius_ratio=0.4)
            txt(s, ix + Inches(0.78), my + Inches(0.10),
                Inches(1.5), Inches(0.3),
                body, size=9, color=INK)
            my += Inches(0.55)
        elif kind == "me":
            shp_round(s, ix + iw - Inches(2.0), my,
                      Inches(1.75), Inches(0.42),
                      BUBBLE_ME, radius_ratio=0.4)
            txt(s, ix + iw - Inches(1.9), my + Inches(0.10),
                Inches(1.55), Inches(0.3),
                body, size=9, color=WHITE)
            my += Inches(0.55)
        elif kind == "place":
            shp_round(s, ix + Inches(0.65), my, Inches(2.0),
                      Inches(0.95), WHITE, line=HAIR, radius_ratio=0.15)
            shp_round(s, ix + Inches(0.75), my + Inches(0.1),
                      Inches(0.55), Inches(0.55),
                      PILL_BG, radius_ratio=0.2)
            txt(s, ix + Inches(0.75), my + Inches(0.18),
                Inches(0.55), Inches(0.4),
                "📍", size=18, color=BRAND, align=PP_ALIGN.CENTER)
            txt(s, ix + Inches(1.4), my + Inches(0.13),
                Inches(1.2), Inches(0.25),
                "OO식당 강남점", size=8, bold=True, color=INK)
            txt(s, ix + Inches(1.4), my + Inches(0.32),
                Inches(1.2), Inches(0.2),
                "Google Maps", size=7, color=MUTED)
            txt(s, ix + Inches(0.75), my + Inches(0.72),
                Inches(2.0), Inches(0.2),
                "탭하여 길찾기 ›", size=7, color=BRAND)
            my += Inches(1.05)
    # input bar
    ib_y = iy + ih - Inches(0.95)
    shp_rect(s, ix, ib_y, iw, Inches(0.005), HAIR)
    shp_oval(s, ix + Inches(0.20), ib_y + Inches(0.13),
             Inches(0.40), Inches(0.40), GRAY_BG)
    txt(s, ix + Inches(0.20), ib_y + Inches(0.18),
        Inches(0.40), Inches(0.3),
        "+", size=14, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    shp_round(s, ix + Inches(0.7), ib_y + Inches(0.14),
              iw - Inches(1.5), Inches(0.40),
              GRAY_BG, radius_ratio=0.5)
    txt(s, ix + Inches(0.85), ib_y + Inches(0.22),
        iw - Inches(1.7), Inches(0.3),
        "메시지", size=10, color=MUTED)
    shp_oval(s, ix + iw - Inches(0.65), ib_y + Inches(0.13),
             Inches(0.40), Inches(0.40), BRAND)
    txt(s, ix + iw - Inches(0.65), ib_y + Inches(0.16),
        Inches(0.40), Inches(0.3),
        "↑", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def draw_place_picker(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    ny = navbar(s, ix, iy, iw, "장소 선택", back=True, action="완료")
    # search
    sy = ny + Inches(0.1)
    shp_round(s, ix + Inches(0.25), sy, iw - Inches(0.5), Inches(0.45),
              GRAY_BG, radius_ratio=0.5)
    txt(s, ix + Inches(0.4), sy + Inches(0.13),
        iw - Inches(0.8), Inches(0.3),
        "🔍  카페, 식당, 주소 검색", size=10, color=MUTED)
    # results
    results = [
        ("OO식당 강남점", "서울 강남구 테헤란로 123", "★ 4.5 · 한식", BRAND),
        ("OO베이커리", "서울 강남구 강남대로 456", "★ 4.7 · 카페·디저트", WARN),
        ("OO이자카야", "서울 강남구 역삼동 789", "★ 4.2 · 일식", BRAND_2),
        ("OO치킨", "서울 강남구 논현동 12", "★ 4.4 · 치킨", DANGER),
    ]
    ry = sy + Inches(0.65)
    for i, (n, addr, meta, c) in enumerate(results):
        y = ry + Inches(i * 0.78)
        shp_round(s, ix + Inches(0.4), y + Inches(0.12),
                  Inches(0.45), Inches(0.45), PILL_BG, radius_ratio=0.3)
        txt(s, ix + Inches(0.4), y + Inches(0.18),
            Inches(0.45), Inches(0.4),
            "📍", size=14, color=c, align=PP_ALIGN.CENTER)
        txt(s, ix + Inches(1.0), y + Inches(0.10),
            iw - Inches(1.3), Inches(0.3),
            n, size=11, bold=True, color=INK)
        txt(s, ix + Inches(1.0), y + Inches(0.36),
            iw - Inches(1.3), Inches(0.25),
            addr, size=8, color=MUTED)
        txt(s, ix + Inches(1.0), y + Inches(0.55),
            iw - Inches(1.3), Inches(0.25),
            meta, size=8, color=c)
        line(s, ix + Inches(1.0), y + Inches(0.78),
             ix + iw, y + Inches(0.78), HAIR, 0.5)


def draw_invite_accept(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    # close
    txt(s, ix + Inches(0.2), iy + Inches(0.15), Inches(0.4), Inches(0.4),
        "✕", size=16, color=MUTED)
    cy = iy + Inches(0.9)
    shp_oval(s, ix + iw / 2 - Inches(0.45), cy,
             Inches(0.9), Inches(0.9), BRAND)
    txt(s, ix + iw / 2 - Inches(0.45), cy + Inches(0.15),
        Inches(0.9), Inches(0.7),
        "🔗", size=32, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ix, cy + Inches(1.2), iw, Inches(0.4),
        "민호님이 초대했어요", size=18, bold=True,
        color=INK, align=PP_ALIGN.CENTER)
    txt(s, ix, cy + Inches(1.6), iw, Inches(0.3),
        "@minho_kim", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # group card
    gy = cy + Inches(2.1)
    shp_round(s, ix + Inches(0.5), gy, iw - Inches(1.0), Inches(0.95),
              SURFACE, line=HAIR, radius_ratio=0.15)
    avatar(s, ix + Inches(0.7), gy + Inches(0.18),
           Inches(0.6), BRAND, "대")
    txt(s, ix + Inches(1.4), gy + Inches(0.20),
        iw - Inches(2), Inches(0.3),
        "대학동기", size=12, bold=True, color=INK)
    txt(s, ix + Inches(1.4), gy + Inches(0.50),
        iw - Inches(2), Inches(0.3),
        "12명의 친구들", size=9, color=MUTED)
    # CTAs
    by = iy + ih - Inches(1.55)
    shp_round(s, ix + Inches(0.25), by, iw - Inches(0.5), Inches(0.5),
              BRAND, radius_ratio=0.5)
    txt(s, ix + Inches(0.25), by + Inches(0.13),
        iw - Inches(0.5), Inches(0.3),
        "수락하고 그룹 참여", size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    shp_round(s, ix + Inches(0.25), by + Inches(0.6),
              iw - Inches(0.5), Inches(0.5), WHITE, line=HAIR,
              radius_ratio=0.5)
    txt(s, ix + Inches(0.25), by + Inches(0.73),
        iw - Inches(0.5), Inches(0.3),
        "거절", size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def draw_profile(s, ix, iy, iw, ih):
    shp_rect(s, ix, iy, iw, ih, APP_BG)
    h = Inches(0.85)
    shp_rect(s, ix, iy, iw, h, WHITE)
    txt(s, ix + Inches(0.3), iy + Inches(0.2), iw, Inches(0.5),
        "내 정보", size=22, bold=True, color=INK)
    # profile card
    py = iy + h + Inches(0.15)
    shp_round(s, ix + Inches(0.25), py, iw - Inches(0.5), Inches(1.6),
              WHITE, line=HAIR, radius_ratio=0.15)
    avatar(s, ix + iw / 2 - Inches(0.45), py + Inches(0.20),
           Inches(0.9), BRAND, "지")
    txt(s, ix, py + Inches(1.15), iw, Inches(0.3),
        "지수", size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    txt(s, ix, py + Inches(1.40), iw, Inches(0.3),
        "@jisoo_28", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    # menu
    items = [
        ("🔔", "기본 알림 설정", "30분 전"),
        ("🌐", "언어", "한국어"),
        ("🔒", "프라이버시", ""),
        ("📩", "친구 초대 보내기", ""),
        ("🚪", "로그아웃", ""),
    ]
    my = py + Inches(1.85)
    for i, (icon, lbl, val) in enumerate(items):
        y = my + Inches(i * 0.50)
        shp_round(s, ix + Inches(0.25), y,
                  iw - Inches(0.5), Inches(0.42),
                  WHITE, line=HAIR, radius_ratio=0.15)
        txt(s, ix + Inches(0.4), y + Inches(0.10),
            Inches(0.3), Inches(0.3),
            icon, size=12, color=BRAND)
        txt(s, ix + Inches(0.8), y + Inches(0.10),
            iw - Inches(2.5), Inches(0.3),
            lbl, size=10, color=INK)
        txt(s, ix + iw - Inches(1.7), y + Inches(0.10),
            Inches(1.4), Inches(0.3),
            val, size=9, color=MUTED, align=PP_ALIGN.RIGHT)
        txt(s, ix + iw - Inches(0.4), y + Inches(0.10),
            Inches(0.2), Inches(0.3),
            "›", size=12, color=MUTED)
    tabbar(s, ix, iy, iw, ih, active=3)


# ============================================================
# Render screens
# ============================================================
SCREENS = [
    (4, 3, "로그인", "LoginScreen", "/auth/login", draw_login,
     ["풀블리드 화이트 + 중앙 로고/CTA",
      "소셜 버튼 3종: Apple(검정) · Google(테두리) · Kakao(노랑)",
      "하단 약관/개인정보 텍스트 (8pt MUTED)"],
     ["버튼 탭 → OAuth 시트 모달",
      "성공 → handle 미설정 시 OnboardingHandleScreen",
      "기존 사용자는 바로 MeetupListScreen"],
     ["로고는 24/24 라운드 사각, 가운데 'M' 글자형",
      "서비스 톤: 친근하고 가벼움. 무거운 일러스트 지양",
      "다크모드 시 배경 INK, 버튼 컬러 그대로 유지"]),
    (5, 4, "온보딩 - 핸들 설정", "OnboardingHandleScreen",
     "/auth/onboarding", draw_onboarding,
     ["원형 아바타 업로드 + 핸들 + 표시이름 2필드",
      "핸들 prefix '@' 고정 표시, 변경 불가 안내",
      "실시간 사용 가능 여부 체크 (✓/✗)"],
     ["핸들 입력 시 디바운스 검증 → 색상 피드백",
      "표시이름은 1~20자 자유, 이모지 허용",
      "시작하기 → 메인 탭으로 이동"],
     ["입력 필드는 GRAY_BG + 포커스 시 BRAND 테두리",
      "오류는 폼 하단 인라인 (DANGER 색)",
      "키보드 위에 CTA 떠있도록 KeyboardAvoidingView"]),
    (6, 5, "약속 목록", "MeetupListScreen", "/meetups", draw_meetup_list,
     ["대형 타이틀 (iOS 스타일)",
      "검색 바 + 세그먼트(다가오는/오늘/지난)",
      "카드: 좌측 컬러 스트라이프로 그룹/약속 구분"],
     ["FAB(+) 탭 → MeetupCreateScreen 모달 push",
      "카드 탭 → MeetupDetailScreen",
      "Pull-to-refresh로 최신 갱신"],
     ["카드 그림자 대신 1px HAIR 테두리 (라이트 톤)",
      "시간 라벨은 카드 색과 동일하게 강조",
      "위치공유/알림 상태는 우측 pill 한 줄"]),
    (7, 6, "약속 상세", "MeetupDetailScreen",
     "/meetups/:id", draw_meetup_detail,
     ["BRAND 색 hero 영역에 약속명·장소",
      "정보 행 3개: 참여자/내 알림/위치공유",
      "지도 미리보기 + 채팅 진입 CTA"],
     ["우측 상단 ••• → 편집/취소/나가기 메뉴",
      "정보 행 탭 → 각 편집 화면",
      "지도 탭 → MeetupMapScreen 풀스크린"],
     ["hero 색은 약속 그룹 색을 따라감 (1회성은 BRAND)",
      "지도 미리보기는 실제 지도 대신 SkeletonStyle",
      "CTA는 sticky로 키보드 위 / 스크롤 무관 항상 노출"]),
    (8, 7, "약속 만들기", "MeetupCreateScreen",
     "/meetups/new", draw_meetup_create,
     ["폼 행 단위 SURFACE 카드 스택",
      "그룹/일시/장소/멤버/위치공유/알림 6필드",
      "헤더 우측 '저장' 비활성→유효시 활성"],
     ["그룹 선택 시 멤버 default 채워짐",
      "장소 탭 → PlacePickerScreen",
      "위치공유 시각: 10/20/30/60분 선택 sheet"],
     ["빈 필드는 placeholder가 라벨처럼 보이도록",
      "필수 미입력 시 행 테두리 DANGER 강조",
      "스크롤 시 헤더 컴팩트 모드 (iOS 패턴)"]),
    (9, 8, "약속 지도", "MeetupMapScreen",
     "/meetups/:id/map", draw_meetup_map,
     ["풀블리드 지도 + 도착지 핀 + 멤버 아바타 핀",
      "하단 시트(드래그 핸들)에 참여자 ETA 리스트",
      "상단 nav 반투명, 지도 위에 떠있음"],
     ["멤버 핀 탭 → 채팅에서 해당 사용자 멘션",
      "시트 드래그 → 풀스크린 리스트 모드",
      "지도 길게 누름 → 임시 핀 (장소 카드 채팅 공유)"],
     ["멤버 색은 프로필 색 일관 (브랜드 토큰화)",
      "위치공유 OFF 멤버는 시트에 'OFF' 회색 처리",
      "줌 레벨은 모든 핀이 보이게 자동 fit"]),
    (10, 9, "그룹 목록", "GroupListScreen", "/groups", draw_group_list,
     ["대형 타이틀 + 우측 상단 '+' (그룹 생성)",
      "카드: 원형 아바타 + 이름 + 멤버수/약속수",
      "카드 우측 chevron(›)"],
     ["카드 탭 → GroupDetailScreen",
      "+ 탭 → GroupCreateScreen",
      "길게 누름 → 컨텍스트 메뉴 (나가기 등)"],
     ["아바타 색은 그룹 ID 해시로 결정 (일관 색상)",
      "그룹이 없으면 일러스트 + '첫 그룹 만들기' 빈 상태",
      "검색 바는 5개 이상부터 노출"]),
    (11, 10, "그룹 상세", "GroupDetailScreen",
     "/groups/:id", draw_group_detail,
     ["상단 카드: 그룹 아바타·이름·멤버수·방장",
      "액션 3버튼: 약속 만들기 / 멤버 / 초대",
      "다가오는 약속 1행 + 그룹 채팅 CTA"],
     ["•••: 그룹 정보 편집 (owner/admin), 나가기",
      "멤버 탭 → GroupMembersScreen (역할 표시)",
      "초대 탭 → GroupInviteScreen (링크/QR)"],
     ["방장 위임은 멤버 화면에서 처리 (이 화면엔 노출 X)",
      "다가오는 약속 0개면 빈 상태 + CTA",
      "스크롤 시 헤더 카드 collapse"]),
    (12, 11, "채팅 목록", "ChatListScreen", "/chats", draw_chat_list,
     ["대형 타이틀 + 채팅 행 리스트",
      "행: 아바타 + 이름 + 마지막 메시지 + 시간 + 뱃지",
      "지난 약속은 'pill: 지난' + MUTED 톤"],
     ["행 탭 → ChatRoomScreen",
      "왼쪽 swipe → 음소거/나가기 액션",
      "검색은 상단 pull-down으로 노출"],
     ["뱃지는 30 이상이면 '30+', 99 이상은 '99+'",
      "지난 약속은 별도 섹션 헤더로 묶을 수 있음 (Phase 2)",
      "마지막 메시지는 1줄 ellipsis"]),
    (13, 12, "채팅방", "ChatRoomScreen",
     "/chats/:room_id", draw_chat_room,
     ["nav 타이틀에 채팅명 + 인원수",
      "위치공유 시작 안내 배너 (PILL_BG)",
      "메시지 버블 + 장소 카드 + 인풋 바"],
     ["인풋 좌측 + → 카메라/이미지/장소 sheet",
      "장소 선택 → PlacePickerScreen 후 장소 카드 전송",
      "메시지 길게 누름 → 편집/삭제/답장"],
     ["내 메시지 우측 BRAND, 상대 메시지 좌측 그레이",
      "장소 카드는 탭 시 외부 지도 앱으로 길찾기",
      "키보드 활성 시 인풋 바가 같이 올라옴"]),
    (14, 13, "장소 선택", "PlacePickerScreen",
     "/place-picker", draw_place_picker,
     ["검색 바 + 결과 리스트 (Google Places)",
      "각 행: 아이콘 + 이름 + 주소 + 별점/카테고리",
      "헤더 우측 '완료' (선택 후 활성)"],
     ["행 탭 → 선택 표시 → 완료로 반환",
      "현재 위치 기준 정렬 (옵션)",
      "최근 선택 장소 상단 노출"],
     ["행 좌측 핀 색은 카테고리별 (한식/카페/일식 등)",
      "검색 0건이면 'OO 검색 결과 없음' 빈 상태",
      "API 비용 위해 디바운스 350ms"]),
    (15, 14, "초대 수락", "InviteAcceptScreen",
     "/invite/:code", draw_invite_accept,
     ["딥링크 진입 시 풀스크린 모달",
      "초대자 정보 + 그룹/친구 카드 + CTA 2개",
      "상단 좌측 ✕ 닫기"],
     ["수락 → 그룹 참여 또는 친구 추가, 해당 화면으로 이동",
      "거절 → 닫기 + 초대 코드 used_count 변경 X",
      "만료/소진 시 'invite expired' 상태 화면"],
     ["로그인 안 된 상태면 LoginScreen 거쳐 다시 진입",
      "이미 멤버면 '이미 참여 중' 상태로 안내",
      "버튼 sticky bottom"]),
    (16, 15, "Profile / Me", "ProfileScreen", "/me", draw_profile,
     ["상단 프로필 카드 (아바타·이름·핸들)",
      "메뉴 행: 알림/언어/프라이버시/초대/로그아웃",
      "각 행 좌측 아이콘 + 우측 값 + chevron"],
     ["행 탭 → 해당 설정 화면",
      "초대 보내기 → 본인 친구 초대 링크 시트",
      "로그아웃 확인 다이얼로그"],
     ["아이콘은 SF Symbols / Material 호환",
      "프라이버시는 위치공유 기본값 + 차단 목록 포함",
      "Phase 2: 다크모드 토글, 알림 채널 세분화"]),
]
for screen in SCREENS:
    page, num, kr, en, route, draw, lay, ix_n, dn = screen
    screen_slide(page, num, kr, en, route, draw, lay, ix_n, dn)


# ─── Final summary ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
shp_rect(s, 0, 0, SW, SH, WHITE)
slide_header(s, 16, "디자인 원칙 요약", "이 화면들 전체에 적용되는 룰")

principles = [
    ("✨", "친근하고 가볍게",
     "친구 사이 앱. 무거운 카드/그림자 지양, 1px HAIR 테두리 + 부드러운 라운딩",
     BRAND),
    ("📍", "위치 공유의 신뢰감",
     "상태(자동 시작/수동 OFF)는 항상 명시적으로 노출. 깜빡 잊을 수 없게",
     SUCCESS),
    ("🎯", "약속 중심 정보 위계",
     "시간·장소·참여자는 항상 같은 위계로. 부가 정보는 pill로 압축",
     WARN),
    ("🔒", "프라이버시는 화면에 보이게",
     "위치공유 ON/OFF, 알림 설정은 한 번에 보이는 곳에. 숨기지 않기",
     BRAND_2),
    ("📱", "iOS 네이티브 패턴 준수",
     "대형 타이틀, 세그먼트 컨트롤, sticky CTA, KeyboardAvoidingView",
     INK),
]
y0 = Inches(1.2)
for i, (icon, t, d, c) in enumerate(principles):
    y = y0 + Inches(i * 1.05)
    shp_round(s, Inches(0.5), y, Inches(12.3), Inches(0.95),
              SURFACE, radius_ratio=0.10)
    shp_round(s, Inches(0.5), y, Inches(0.10), Inches(0.95), c,
              radius_ratio=0.3)
    txt(s, Inches(0.8), y + Inches(0.25), Inches(0.7), Inches(0.5),
        icon, size=24, color=c)
    txt(s, Inches(1.7), y + Inches(0.18), Inches(10.5), Inches(0.4),
        t, size=15, bold=True, color=INK)
    txt(s, Inches(1.7), y + Inches(0.55), Inches(10.5), Inches(0.4),
        d, size=11, color=MUTED)
slide_footer(s, 16, TOTAL)


out = "d:/Workspace/CPKWorks/MeetPod/docs/MeetPod_화면설계.pptx"
prs.save(out)
