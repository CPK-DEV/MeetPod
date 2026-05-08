"""MeetPod 사용 시나리오 14개 PPT (flow chart 포함)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x14, 0x2A, 0x5C)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PINK = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

TOTAL = 17  # cover + persona + 14 stories + priority


def rect(slide, x, y, w, h, fill, line=None, shadow=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    if not shadow: s.shadow.inherit = False
    return s


def round_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, t, size=14, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, vcenter=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if vcenter:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return tb


def arrow(slide, x1, y1, x2, y2, color=NAVY):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(2.5)
    # arrow head
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = c.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return c


def header(slide, title, subtitle=None, color=NAVY):
    rect(slide, 0, 0, SW, Inches(0.9), color)
    text(slide, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.5),
         title, size=22, bold=True, color=WHITE)
    if subtitle:
        text(slide, Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.3),
             subtitle, size=11, color=LIGHT)


def footer(slide, page):
    text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
         "MeetPod — 사용자 스토리", size=9, color=MUTED)
    text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
         f"{page} / {TOTAL}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def persona_badge(slide, x, y, name, age_job, color):
    """원형 아바타 + 이름표"""
    oval(slide, x, y, Inches(0.7), Inches(0.7), color)
    text(slide, x, y + Inches(0.13), Inches(0.7), Inches(0.5),
         name[0], size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.85), y + Inches(0.05), Inches(3.5), Inches(0.35),
         name, size=14, bold=True, color=DARK)
    text(slide, x + Inches(0.85), y + Inches(0.4), Inches(3.5), Inches(0.3),
         age_job, size=10, color=MUTED)


def situation_box(slide, x, y, w, text_str):
    round_rect(slide, x, y, w, Inches(0.55), LIGHT)
    text(slide, x + Inches(0.2), y + Inches(0.13), w - Inches(0.4), Inches(0.3),
         f"💡 상황: {text_str}", size=11, color=DARK)


def value_box(slide, x, y, w, text_str):
    round_rect(slide, x, y, w, Inches(0.6), GREEN)
    text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3),
         f"✓ 가치: {text_str}", size=12, bold=True, color=WHITE)


def flow_horizontal(slide, x0, y0, steps, box_w=Inches(2.0), box_h=Inches(1.3),
                    gap=Inches(0.35), color=ACCENT):
    """가로 flow: 박스 → 화살표 → 박스 ..."""
    n = len(steps)
    for i, step in enumerate(steps):
        x = x0 + (box_w + gap) * i
        round_rect(slide, x, y0, box_w, box_h, WHITE, line=color)
        # number circle
        oval(slide, x + Inches(0.05), y0 + Inches(0.05),
             Inches(0.32), Inches(0.32), color)
        text(slide, x + Inches(0.05), y0 + Inches(0.07),
             Inches(0.32), Inches(0.28),
             str(i + 1), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.15), y0 + Inches(0.45),
             box_w - Inches(0.3), box_h - Inches(0.5),
             step, size=10, color=DARK, align=PP_ALIGN.CENTER, vcenter=True)
        if i < n - 1:
            ax1 = x + box_w
            ax2 = x + box_w + gap
            ay = y0 + box_h / 2
            arrow(slide, ax1, ay, ax2, ay, color=color)


def flow_vertical(slide, x0, y0, steps, box_w=Inches(8.5), box_h=Inches(0.65),
                  gap=Inches(0.18), color=ACCENT):
    n = len(steps)
    for i, step in enumerate(steps):
        y = y0 + (box_h + gap) * i
        round_rect(slide, x0, y, box_w, box_h, WHITE, line=color)
        oval(slide, x0 + Inches(0.1), y + Inches(0.13),
             Inches(0.4), Inches(0.4), color)
        text(slide, x0 + Inches(0.1), y + Inches(0.16),
             Inches(0.4), Inches(0.35),
             str(i + 1), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, x0 + Inches(0.65), y + Inches(0.05),
             box_w - Inches(0.8), box_h - Inches(0.1),
             step, size=11, color=DARK, vcenter=True)
        if i < n - 1:
            cx = x0 + Inches(0.3)
            arrow(slide, cx, y + box_h, cx, y + box_h + gap, color=color)


def flow_grid(slide, x0, y0, steps, cols=4, box_w=Inches(2.85),
              box_h=Inches(1.1), gap_x=Inches(0.2), gap_y=Inches(0.3),
              color=ACCENT):
    """4열 grid, 행 간 ↓ 화살표"""
    for i, step in enumerate(steps):
        col = i % cols
        row = i // cols
        x = x0 + (box_w + gap_x) * col
        y = y0 + (box_h + gap_y) * row
        round_rect(slide, x, y, box_w, box_h, WHITE, line=color)
        oval(slide, x + Inches(0.05), y + Inches(0.05),
             Inches(0.32), Inches(0.32), color)
        text(slide, x + Inches(0.05), y + Inches(0.07),
             Inches(0.32), Inches(0.28),
             str(i + 1), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.15), y + Inches(0.4),
             box_w - Inches(0.3), box_h - Inches(0.45),
             step, size=9, color=DARK, align=PP_ALIGN.CENTER, vcenter=True)
        # arrow within row
        if col < cols - 1 and i + 1 < len(steps):
            ax1 = x + box_w
            ax2 = x + box_w + gap_x
            ay = y + box_h / 2
            arrow(slide, ax1, ay, ax2, ay, color=color)
        # downward arrow at end of row
        if col == cols - 1 and i + 1 < len(steps):
            ay1 = y + box_h
            ay2 = y + box_h + gap_y
            ax = x + box_w / 2
            arrow(slide, ax, ay1, ax, ay2, color=color)


def story_slide(page, story_no, title, persona_name, persona_meta,
                persona_color, situation, steps, value, color=ACCENT,
                layout="grid"):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    header(s, f"스토리 {story_no} — {title}",
           f"페르소나 시나리오 #{story_no}", color=color)
    persona_badge(s, Inches(0.5), Inches(1.1), persona_name,
                  persona_meta, persona_color)
    situation_box(s, Inches(5.3), Inches(1.18), Inches(7.5), situation)
    if layout == "grid":
        flow_grid(s, Inches(0.5), Inches(2.3), steps, color=color)
    elif layout == "vertical":
        flow_vertical(s, Inches(2.4), Inches(2.3), steps, color=color)
    elif layout == "horizontal":
        flow_horizontal(s, Inches(0.5), Inches(2.3), steps, color=color)
    value_box(s, Inches(0.5), Inches(6.4), Inches(12.3), value)
    footer(s, page)
    return s


# ─────────────────────────────────────────────────────────
# Slide 1: 표지
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.0), SW, Inches(0.04), ACCENT)
text(s, Inches(0.8), Inches(1.8), Inches(12), Inches(1.0),
     "MeetPod", size=60, bold=True, color=WHITE)
text(s, Inches(0.8), Inches(2.8), Inches(12), Inches(0.5),
     "사용자 스토리 14선", size=28, bold=True, color=ACCENT)
text(s, Inches(0.8), Inches(3.5), Inches(12), Inches(0.5),
     "페르소나로 본 MeetPod 사용 시나리오와 흐름",
     size=18, color=LIGHT)
# decorative dots
for i, c in enumerate([ACCENT, GREEN, ORANGE, PURPLE, PINK, RED]):
    oval(s, Inches(0.8 + i * 0.6), Inches(5.5), Inches(0.4), Inches(0.4), c)
text(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
     "2026.05.03  ·  CPKWorks", size=12, color=MUTED)

# ─────────────────────────────────────────────────────────
# Slide 2: 페르소나 6명 개요
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, "페르소나 6인", "스토리에 등장하는 사용자")
personas = [
    ("지수", "28 · 직장인", "퇴근 후 친구 모임", ACCENT),
    ("민호", "24 · 대학생", "동기 모임, 즉흥 약속", GREEN),
    ("서연", "32 · 워킹맘", "육아맘 친구 조율", ORANGE),
    ("태훈", "35 · 동호회장", "등산 동호회 12명", PURPLE),
    ("혜진", "22 · 신입생", "새 친구 사귀는 중", PINK),
    ("준영", "40 · 자영업", "사장님 비정기 회식", RED),
]
x0, y0 = Inches(0.6), Inches(1.4)
cw, ch, gx, gy = Inches(4.0), Inches(2.7), Inches(0.18), Inches(0.25)
for i, (name, meta, ctx, color) in enumerate(personas):
    col, row = i % 3, i // 3
    x = x0 + (cw + gx) * col
    y = y0 + (ch + gy) * row
    round_rect(s, x, y, cw, ch, LIGHT)
    rect(s, x, y, cw, Inches(0.06), color)
    oval(s, x + Inches(1.4), y + Inches(0.4),
         Inches(1.2), Inches(1.2), color)
    text(s, x + Inches(1.4), y + Inches(0.65),
         Inches(1.2), Inches(0.7),
         name[0], size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(1.7), cw, Inches(0.4),
         name, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(2.05), cw, Inches(0.3),
         meta, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(2.35), cw, Inches(0.3),
         ctx, size=11, color=DARK, align=PP_ALIGN.CENTER)
footer(s, 2)

# ─────────────────────────────────────────────────────────
# Story 1
story_slide(3, 1, "퇴근 후 갑작스런 약속 (단발성)", "지수", "28 · 직장인",
            ACCENT,
            "금요일 오후 5시. 친구 둘과 갑자기 저녁 약속이 잡혔다.",
            ["카톡으로 '7시 강남역' 제안 받음",
             "MeetPod '+ 새 약속' 탭",
             "1회성 약속 선택",
             "친구 2명 선택",
             "장소: OO식당 (Google Maps)",
             "19:00~21:00, 위치공유 20분 전",
             "본인 알림 30분 전 추가",
             "저장 → 친구에게 푸시",
             "18:40 자동 위치공유 시작",
             "18:55 지도에 친구 핀 확인",
             "메시지 없이 바로 입장",
             "21:00 자동 종료, 채팅 아카이브"],
            "'어디야?' 카톡 없이 도착 확인. 종료 시 위치공유 자동 OFF.",
            color=ACCENT)

# Story 2
story_slide(4, 2, "대학 동기 동창회 (그룹 + 약속)", "민호", "24 · 대학생",
            GREEN,
            "학과 동기 15명 단체. 1년 만의 모임.",
            ["'대학동기' 그룹 생성 (방장)",
             "초대 링크 → 학과 단톡방 공유",
             "12명 가입 완료",
             "'5/15 동창회' 약속 생성",
             "12명 default, 출장자 1명 해제",
             "장소: 홍대 OO술집 (30분 전 공유)",
             "약속 채팅방 자동 생성",
             "당일 더 좋은 가게 발견",
             "장소 카드로 채팅 공유",
             "민호가 약속 편집 → 변경 푸시",
             "지도로 도착 현황 한눈에",
             "종료 후 그룹 채팅에서 '다음은?'"],
            "큰 그룹 조율 + 장소 변경 푸시 + 그룹/약속 채팅 분리.",
            color=GREEN)

# Story 3
story_slide(5, 3, "장소가 공사 중 (실시간 변경)", "서연", "32 · 워킹맘",
            ORANGE,
            "육아맘 친구 3명과 평일 낮 카페 약속.",
            ["일주일 전 약속 생성",
             "당일 11:30 위치공유 자동 ON",
             "한 친구 도착 → 공사 중 발견",
             "채팅에 상황 + 옆 카페 장소 카드",
             "서연이 약속 편집 → 새 카페로",
             "다른 멤버 푸시 + 핀 자동 업데이트",
             "모두 새 카페로 이동",
             "위치공유 덕에 '5분 뒤 도착' 채팅 불필요"],
            "실시간 장소 변경, 채팅 장소 카드, Realtime 핀 업데이트.",
            color=ORANGE)

# Story 4
story_slide(6, 4, "등산 동호회 산행 (위치 공유의 진가)", "태훈", "35 · 동호회장",
            PURPLE,
            "12명 등산 동호회. 토요일 새벽 산행.",
            ["'OO산악회' 그룹 운영",
             "부회장에게 admin 위임",
             "토 06:00~14:00 약속",
             "위치공유 60분 전으로 변경",
             "새벽 5시 멤버 자동 트래킹",
             "06:00 11명 도착, 1명 30분 거리",
             "채팅에 '먼저 출발' 공지",
             "산행 중 한 명 컨디션 난조",
             "조기 하산 채팅 알림",
             "지도에서 위치 확인 가능",
             "14:00 종료, 채팅 아카이브",
             "그룹 채팅에서 사진 공유 + 다음 약속"],
            "다인원 위치 가시성 + 안전 + owner/admin 권한 분리.",
            color=PURPLE)

# Story 5
story_slide(7, 5, "친구가 길을 잃음 (안전)", "지수", "28 · 직장인",
            RED,
            "처음 가는 동네에서 만나기로 함.",
            ["약속 시간 10분 전",
             "지수 도착, 친구 미도착",
             "카톡 대신 지도 확인",
             "친구가 반대편 골목에 핀",
             "채팅에 '큰길로 나와 오른쪽'",
             "본인 장소 다시 카드로 공유",
             "친구 5분 후 도착",
             "전화 없이 해결"],
            "전화 안 하고도 길 안내. 일상의 작은 마찰 제거.",
            color=RED, layout="horizontal")

# Story 6
story_slide(8, 6, "신입생 첫 가입 (온보딩)", "혜진", "22 · 신입생",
            PINK,
            "대학 신입생. 친구가 MeetPod 추천.",
            ["친구가 카톡으로 초대 링크 전송",
             "링크 탭 → 앱스토어 → 설치",
             "Kakao 로그인",
             "@hyejin22 핸들 입력",
             "푸시 권한 허용",
             "친구 신청 자동 수락",
             "신입생 환영회 약속 자동 초대 확인",
             "참여 자동 + 알림 1시간 전 설정",
             "당일 위치공유 권한 '항상 허용'",
             "약속 시간에 자동 트래킹 시작"],
            "가입~첫 약속 참여까지 마찰 최소화. 막히는 곳 없음.",
            color=PINK)

# Story 7
story_slide(9, 7, "사장님 비정기 회식 (다중 그룹)", "준영", "40 · 자영업",
            ORANGE,
            "지역상인회 + 동종업종 사장님 그룹 운영.",
            ["그룹: 상인회(25명) + 치킨집(8명)",
             "상인회 일부 8명만 회식 약속",
             "멤버 picker 25명 → 17명 해제",
             "약속 채팅방 별도 생성",
             "그룹 상시 채팅 조용히 유지",
             "다른 날 치킨집 회의 약속",
             "두 약속 홈에 시간 순 정렬",
             "각 약속 알림 분리"],
            "그룹 부분 모임 + 약속별 채팅으로 단톡 보호.",
            color=ORANGE, layout="horizontal")

# Story 8
story_slide(10, 8, "약속 취소 / 노쇼 (예외)", "지수", "28 · 직장인",
            RED,
            "약속 1시간 전, 한 명이 갑자기 못 옴.",
            ["지수가 채팅에 '못 가게 됐어'",
             "만든 친구가 약속 편집",
             "멤버에서 지수 제거",
             "지수 약속 목록에서 사라짐",
             "지수 위치공유 자동 시작 안 함",
             "[전체 취소 시] 상태 'cancelled'",
             "참여자에게 취소 푸시",
             "채팅 즉시 아카이브, 위치공유 X"],
            "부분 이탈 vs 전체 취소 분리. 위치공유 안전 처리.",
            color=RED, layout="horizontal")

# Story 9
story_slide(11, 9, "위치 공유 거부 (프라이버시)", "서연", "32 · 워킹맘",
            GREEN,
            "약속에는 가지만 위치 공유는 하기 싫음.",
            ["약속 상세 화면 진입",
             "본인 토글: '내 위치 공유 안 함' OFF",
             "본인 위치공유 시작 안 함",
             "다른 멤버 지도에 본인 핀 X",
             "다른 사람 위치는 수신 가능",
             "약속별 설정, 다음 약속 영향 X"],
            "위치공유는 강제 아닌 선택. 친구 신뢰 기반.",
            color=GREEN, layout="horizontal")

# Story 10
story_slide(12, 10, "그룹 방장 위임 / 탈퇴 (관리)", "태훈", "35 · 동호회장",
            PURPLE,
            "동호회를 다른 사람에게 넘기고 탈퇴.",
            ["멤버 화면 → 부회장 선택",
             "'방장 위임' 탭 + 확인",
             "부회장 owner 승격, 태훈 admin 강등",
             "본인 '그룹 나가기' 탭",
             "그룹에서 제거",
             "진행 중 약속에서 자동 제외",
             "새 owner가 그룹 운영 지속"],
            "방장 권한 안전 이전. owner 부재 방지.",
            color=PURPLE, layout="horizontal")

# Story 11
story_slide(13, 11, "약속 종료 후 채팅 (사후 정리)", "민호", "24 · 대학생",
            ACCENT,
            "약속 끝났지만 사진 공유 등 사후 대화 필요.",
            ["종료 시각 도달",
             "채팅방 자동 아카이브",
             "'지난 약속' 섹션으로 이동",
             "진입/송수신 가능",
             "민호가 사진 공유",
             "다른 참여자 수신",
             "위치공유는 종료 상태 유지",
             "24시간 후 위치 핑 자동 삭제"],
            "사후 정리 가능 + 화면은 깔끔하게 분리.",
            color=ACCENT, layout="horizontal")

# Story 12
story_slide(14, 12, "메시지 편집 / 삭제", "혜진", "22 · 신입생",
            PINK,
            "채팅에 잘못된 정보 보냄.",
            ["본인 메시지 길게 누름",
             "메뉴: 편집 / 삭제",
             "[편집] 텍스트 수정",
             "edited_at 표시 노출",
             "[삭제] soft delete",
             "다른 멤버에게 '삭제된 메시지' 표시"],
            "실수 복구 + 투명한 편집 표시.",
            color=PINK, layout="horizontal")

# ─────────────────────────────────────────────────────────
# Story 13: 푸시 알림 종류 (table)
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, "스토리 13 — 푸시 알림 다양한 케이스",
       "사용자에게 도달하는 알림 7가지", color=ORANGE)
situation_box(s, Inches(0.5), Inches(1.1), Inches(12.3),
              "알림은 너무 많지도 적지도 않게. 각 케이스별 문구 가이드.")
notifs = [
    ("📨", "새 약속 초대",
     "'지수님이 ′강남역 저녁′에 초대했어요 (오늘 19:00)'"),
    ("✏️", "약속 정보 변경",
     "'약속 장소가 ′OO식당′으로 변경되었어요'"),
    ("❌", "약속 취소",
     "'오늘 19:00 약속이 취소되었어요'"),
    ("⏰", "본인 설정 알림",
     "'약속 30분 전이에요 — 강남역 저녁'"),
    ("📡", "위치공유 시작 임박",
     "'20분 뒤 위치 공유가 시작돼요. 끄려면 약속 화면에서 토글'"),
    ("👥", "그룹 초대",
     "'민호님이 ′대학동기′ 그룹에 초대했어요'"),
    ("💬", "새 채팅 메시지",
     "(포그라운드 외) 사용자별 ON/OFF 가능"),
]
y0 = Inches(2.0)
rh = Inches(0.65)
for i, (icon, kind, msg) in enumerate(notifs):
    y = y0 + (rh + Inches(0.08)) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    round_rect(s, Inches(0.5), y, Inches(12.3), rh, bg, line=ORANGE)
    text(s, Inches(0.7), y + Inches(0.16), Inches(0.6), Inches(0.4),
         icon, size=20)
    text(s, Inches(1.4), y + Inches(0.18), Inches(2.8), Inches(0.4),
         kind, size=13, bold=True, color=DARK)
    text(s, Inches(4.3), y + Inches(0.2), Inches(8.3), Inches(0.4),
         msg, size=11, color=MUTED)
footer(s, 15)

# Story 14
story_slide(16, 14, "해외 여행 / 시차 (엣지)", "지수", "28 · 직장인",
            MUTED,
            "친구가 일본 여행 중인데 한국 약속을 잡아둠.",
            ["약속은 한국 시간 (KST) 저장",
             "일본 친구 모바일은 JST로 표시",
             "위치공유 시작 시각 정확히 동기화",
             "친구가 한국에 없으면 핀이 의미 X",
             "다른 멤버에게 '○○님 1500km 떨어짐' 안내",
             "거리/시차 안내는 Phase 2 후속"],
            "Timezone 정확성 + 원거리 멤버 부드러운 안내.",
            color=MUTED, layout="horizontal")

# ─────────────────────────────────────────────────────────
# Slide 17: 우선순위 매핑
# ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, "스토리 → MVP Phase 매핑", "구현 순서로 본 우선순위")

pri = [
    ("Phase 1: 인증 + 그룹/친구",
     ["6 신입생 온보딩", "7 다중 그룹", "10 방장 위임"], GREEN),
    ("Phase 2: 약속 + 알림",
     ["1 단발 약속(생성 부분)", "8 약속 취소", "13 푸시 알림"], ACCENT),
    ("Phase 3: 채팅",
     ["2 동창회(채팅)", "3 장소 변경", "11 사후 정리", "12 편집/삭제"], ORANGE),
    ("Phase 4: 위치 공유",
     ["1 단발 약속(공유 부분)", "4 산행 안전", "5 길 안내", "9 위치 거부"], PURPLE),
    ("Phase 5+: 후속",
     ["14 시차 안내"], MUTED),
]
x0, y0 = Inches(0.5), Inches(1.4)
cw = Inches(2.5)
gap = Inches(0.07)
for i, (phase, stories, color) in enumerate(pri):
    x = x0 + (cw + gap) * i
    rect(s, x, y0, cw, Inches(0.7), color)
    text(s, x + Inches(0.15), y0 + Inches(0.18), cw - Inches(0.3),
         Inches(0.4),
         phase, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, y0 + Inches(0.7), cw, Inches(4.5), LIGHT)
    sy = y0 + Inches(0.9)
    for st in stories:
        round_rect(s, x + Inches(0.15), sy, cw - Inches(0.3),
                   Inches(0.55), WHITE, line=color)
        text(s, x + Inches(0.25), sy + Inches(0.13),
             cw - Inches(0.5), Inches(0.35),
             st, size=10, color=DARK)
        sy += Inches(0.65)

round_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), LIGHT)
text(s, Inches(0.7), Inches(6.48), Inches(12), Inches(0.4),
     "각 Phase는 독립 배포 가능. 스토리 단위로 QA 테스트 케이스 도출.",
     size=12, color=DARK)
footer(s, 17)

out = "d:/Workspace/CPKWorks/MeetPod/docs/MeetPod_사용자_스토리.pptx"
prs.save(out)
