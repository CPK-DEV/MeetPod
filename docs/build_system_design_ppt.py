"""MeetPod 시스템 설계 PPT 빌드."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0F, 0x1B, 0x3D)
BRAND = RGBColor(0x5B, 0x7C, 0xFA)
HERMES = RGBColor(0xFF, 0x69, 0x00)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARN = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_2 = RGBColor(0x33, 0x40, 0x5C)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
HAIR = RGBColor(0xE2, 0xE8, 0xF0)
SURFACE = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
FONT = "맑은 고딕"


def _i(v):
    """모든 EMU 값을 int로 강제 (float은 PowerPoint에서 invalid)."""
    return int(v) if v is not None else None


def rect(s, x, y, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            _i(x), _i(y), _i(w), _i(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def rrect(s, x, y, w, h, fill, line=None, r=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            _i(x), _i(y), _i(w), _i(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if r is not None: sh.adjustments[0] = r
    return sh


def oval(s, x, y, w, h, fill, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            _i(x), _i(y), _i(w), _i(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, t, size=11, bold=False, color=INK,
        align=PP_ALIGN.LEFT, vcenter=False):
    tb = s.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if vcenter: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT
    return tb


def _add_dash(ln):
    """prstDash를 schema 순서에 맞게 삽입 (fill 뒤, headEnd/tailEnd 앞)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    # 기존 prstDash 제거
    for existing in ln.findall(qn('a:prstDash')):
        ln.remove(existing)
    prst = etree.SubElement(ln, qn('a:prstDash'))
    prst.set('val', 'dash')
    # tailEnd / headEnd 보다 앞으로 이동
    head = ln.find(qn('a:headEnd'))
    tail = ln.find(qn('a:tailEnd'))
    anchor = head if head is not None else tail
    if anchor is not None:
        anchor.addprevious(prst)


def _add_tail_arrow(ln):
    from pptx.oxml.ns import qn
    from lxml import etree
    for existing in ln.findall(qn('a:tailEnd')):
        ln.remove(existing)
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')


def _safe_xy(x1, y1, x2, y2):
    """가로/세로 connector cy=0/cx=0 회피용 1 EMU 오프셋."""
    x1, y1, x2, y2 = int(x1), int(y1), int(x2), int(y2)
    if x1 == x2:
        x2 += 1
    if y1 == y2:
        y2 += 1
    return x1, y1, x2, y2


def arrow(s, x1, y1, x2, y2, color=NAVY, weight=2.0, dashed=False):
    x1, y1, x2, y2 = _safe_xy(x1, y1, x2, y2)
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    ln = c.line._get_or_add_ln()
    if dashed:
        _add_dash(ln)
    _add_tail_arrow(ln)
    return c


def line(s, x1, y1, x2, y2, color=HAIR, weight=0.75, dashed=False):
    x1, y1, x2, y2 = _safe_xy(x1, y1, x2, y2)
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(weight)
    if dashed:
        ln = c.line._get_or_add_ln()
        _add_dash(ln)
    return c


def header(s, num, title, subtitle):
    rect(s, 0, 0, SW, Inches(0.85), NAVY)
    rect(s, 0, Inches(0.85), SW, Inches(0.04), HERMES)
    txt(s, Inches(0.4), Inches(0.13), Inches(0.8), Inches(0.5),
        f"{num:02d}", size=22, bold=True, color=HERMES)
    txt(s, Inches(1.2), Inches(0.13), Inches(11), Inches(0.4),
        title, size=20, bold=True, color=WHITE)
    txt(s, Inches(1.2), Inches(0.50), Inches(11), Inches(0.3),
        subtitle, size=11, color=HAIR)


def footer(s, page, total):
    txt(s, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
        "MeetPod — 시스템 설계", size=9, color=MUTED)
    txt(s, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
        f"{page} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def block(s, x, y, w, h, title, subtitle="", color=BRAND):
    rrect(s, x, y, w, h, WHITE, line=color, r=0.10)
    rect(s, x, y, w, Inches(0.40), color)
    txt(s, x, y + Inches(0.07), w, Inches(0.3),
        title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(s, x + Inches(0.15), y + Inches(0.55),
            w - Inches(0.3), h - Inches(0.6),
            subtitle, size=9, color=INK_2)


TOTAL = 18

# ─────────────────────────────────────────────
# 1. Cover
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.0), SW, Inches(0.05), HERMES)
rect(s, 0, Inches(3.10), SW, Inches(0.02), BRAND)
txt(s, Inches(0.8), Inches(1.6), Inches(12), Inches(1.2),
    "MeetPod", size=64, bold=True, color=WHITE)
txt(s, Inches(0.8), Inches(2.6), Inches(12), Inches(0.5),
    "System Design", size=28, bold=True, color=HERMES)
txt(s, Inches(0.8), Inches(3.4), Inches(12), Inches(0.5),
    "아키텍처 · 시퀀스 · 모듈 · API · 운영 · 위험",
    size=16, color=HAIR)
txt(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
    "2026.05.03  ·  CPKWorks", size=11, color=MUTED)

# ─────────────────────────────────────────────
# 2. Architecture
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 1, "Architecture", "전체 컴포넌트 구성")

# Mobile
block(s, Inches(0.5), Inches(1.3), Inches(3.3), Inches(2.5),
      "Mobile (Expo RN)",
      "UI / Screens\nNavigation\nZustand stores\napi · lib · store",
      BRAND)
# FastAPI
block(s, Inches(5.0), Inches(1.3), Inches(3.3), Inches(2.5),
      "Backend (FastAPI)",
      "Vercel Serverless\nrouters / services\ndependencies / utils\nJWT 검증",
      SUCCESS)
# Supabase
block(s, Inches(9.5), Inches(1.3), Inches(3.4), Inches(2.5),
      "Supabase",
      "Auth (OAuth)\nPostgres + RLS\nRealtime · Storage\npg_cron + Edge Fn",
      HERMES)

# arrows
arrow(s, Inches(3.8), Inches(2.2), Inches(5.0), Inches(2.2), NAVY)
txt(s, Inches(3.8), Inches(1.95), Inches(1.2), Inches(0.25),
    "HTTPS REST", size=8, color=MUTED, align=PP_ALIGN.CENTER)
arrow(s, Inches(8.3), Inches(2.2), Inches(9.5), Inches(2.2), NAVY)
txt(s, Inches(8.3), Inches(1.95), Inches(1.2), Inches(0.25),
    "supabase-py", size=8, color=MUTED, align=PP_ALIGN.CENTER)
# direct mobile→supabase
arrow(s, Inches(2.2), Inches(3.8), Inches(11.0), Inches(3.8),
      HERMES, dashed=True)
txt(s, Inches(5.5), Inches(3.85), Inches(3.0), Inches(0.25),
    "WebSocket (Realtime) · OAuth · Storage upload",
    size=8, color=HERMES, align=PP_ALIGN.CENTER, bold=True)

# bottom: Expo Push
block(s, Inches(5.0), Inches(5.2), Inches(3.3), Inches(1.2),
      "Expo Push API",
      "Edge Function이 reminder/event 발송", PURPLE)
arrow(s, Inches(11.0), Inches(3.7), Inches(8.3), Inches(5.6), PURPLE)
txt(s, Inches(8.5), Inches(4.5), Inches(2.5), Inches(0.3),
    "Edge Fn → Push", size=9, color=PURPLE)

footer(s, 2, TOTAL)

# ─────────────────────────────────────────────
# 3. Traffic Routing
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 2, "Traffic Routing", "어떤 요청이 어디로 가는가")

routing = [
    ("회원가입/로그인", "Mobile ↔ Supabase Auth", "OAuth 직접", BRAND),
    ("프로필 부트스트랩", "Mobile → FastAPI", "handle 충돌 검증·트랜잭션", SUCCESS),
    ("그룹/약속 CRUD", "Mobile → FastAPI", "권한·트랜잭션·트리거", SUCCESS),
    ("초대 코드 발급/소비", "Mobile → FastAPI", "충돌·만료·max_uses 원자성", SUCCESS),
    ("채팅 송신", "Mobile → FastAPI", "권한·archived 체크·push fanout", SUCCESS),
    ("채팅 수신", "Mobile ← Supabase Realtime", "WS 직결, FastAPI 우회", HERMES),
    ("위치 핑 송신", "Mobile → Supabase 직접 INSERT", "10초 주기, latency 민감", HERMES),
    ("위치 핑 수신", "Mobile ← Supabase Realtime", "동일", HERMES),
    ("푸시 발송", "Edge Function → Expo Push", "DB 가까운 곳에서 fanout", PURPLE),
    ("이미지 업로드", "Mobile → Supabase Storage", "직접 (signed URL/RLS)", HERMES),
]
hy = Inches(1.2)
rect(s, Inches(0.5), hy, Inches(12.3), Inches(0.45), NAVY)
for i, (h, c) in enumerate([("동작", 3.5), ("경로", 4.5), ("이유", 4.3)]):
    cx = [Inches(0.5), Inches(4.0), Inches(8.5)][i]
    txt(s, cx + Inches(0.2), hy + Inches(0.10), Inches(c), Inches(0.3),
        h, size=11, bold=True, color=WHITE)
for i, (a, b, c, color) in enumerate(routing):
    y = hy + Inches(0.55) + Inches(i * 0.50)
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.46), bg)
    rect(s, Inches(0.5), y, Inches(0.10), Inches(0.46), color)
    txt(s, Inches(0.7), y + Inches(0.13), Inches(3.3), Inches(0.3),
        a, size=10, bold=True, color=INK)
    txt(s, Inches(4.0), y + Inches(0.13), Inches(4.5), Inches(0.3),
        b, size=10, color=INK_2)
    txt(s, Inches(8.5), y + Inches(0.13), Inches(4.3), Inches(0.3),
        c, size=10, color=MUTED)
footer(s, 3, TOTAL)

# ─────────────────────────────────────────────
# 4. Auth Flow
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 3, "Auth Flow", "Supabase JWT 단일 진실 공급원")

steps = [
    ("Mobile",
     "Supabase JS SDK로 OAuth (Google/Apple/Kakao)",
     BRAND),
    ("Supabase Auth",
     "JWT 발급: access_token + refresh_token",
     HERMES),
    ("Mobile → FastAPI",
     "Authorization: Bearer <access_token>",
     SUCCESS),
    ("FastAPI",
     "SUPABASE_JWT_SECRET로 서명 검증\nsub 클레임 → user_id 추출",
     PURPLE),
    ("FastAPI",
     "자체 JWT 발급 X (Supabase가 단일 SoT)",
     INK),
]
y0 = Inches(1.3)
for i, (who, what, color) in enumerate(steps):
    y = y0 + Inches(i * 0.85)
    oval(s, Inches(0.6), y, Inches(0.55), Inches(0.55), color)
    txt(s, Inches(0.6), y + Inches(0.1), Inches(0.55), Inches(0.4),
        str(i + 1), size=16, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    rrect(s, Inches(1.4), y, Inches(11.3), Inches(0.7),
          SURFACE, line=color, r=0.15)
    txt(s, Inches(1.6), y + Inches(0.07),
        Inches(2.5), Inches(0.3),
        who, size=10, bold=True, color=color)
    txt(s, Inches(4.2), y + Inches(0.10),
        Inches(8.4), Inches(0.5),
        what, size=11, color=INK)

# Kakao note
y = y0 + Inches(5 * 0.85) + Inches(0.1)
rrect(s, Inches(0.6), y, Inches(12.1), Inches(0.7),
      RGBColor(0xFE, 0xF3, 0xC7), line=WARN, r=0.1)
txt(s, Inches(0.8), y + Inches(0.08),
    Inches(12), Inches(0.3),
    "⚠ Kakao 예외", size=10, bold=True, color=WARN)
txt(s, Inches(0.8), y + Inches(0.32),
    Inches(11.8), Inches(0.4),
    "Supabase가 Kakao OIDC native 미지원 시: (1) Custom OIDC provider, "
    "(2) FastAPI가 Kakao 토큰 검증 후 Supabase Admin API로 user 발급",
    size=9, color=INK_2)
footer(s, 4, TOTAL)


# ─────────────────────────────────────────────
# Sequence diagram helper
# ─────────────────────────────────────────────
def seq_diagram(slide, lanes, steps, x0=Inches(0.6), y0=Inches(1.15),
                width=Inches(12.1), step_h=None):
    """
    lanes: [(label, color), ...]
    steps: [(from_idx, to_idx, label, color?), ...]  — to_idx=from for self
    """
    n = len(lanes)
    lane_w = width / n
    # 자동 step 높이 계산: footer(7.0)까지 들어가게
    n_steps = len(steps)
    avail_h = Inches(7.0) - (y0 + Inches(0.55))
    if step_h is None:
        max_h = Inches(0.40)
        step_h = min(max_h, avail_h / max(n_steps + 1, 1))
    # lane headers
    for i, (lbl, c) in enumerate(lanes):
        cx = x0 + lane_w * i
        rrect(slide, cx + Inches(0.05), y0,
              lane_w - Inches(0.1), Inches(0.45), c, r=0.2)
        txt(slide, cx, y0 + Inches(0.10), lane_w, Inches(0.3),
            lbl, size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
    # lifelines
    line_top = y0 + Inches(0.55)
    line_bot = line_top + step_h * (n_steps + 1)
    for i in range(n):
        cx = x0 + lane_w * i + lane_w / 2
        line(slide, cx, line_top, cx, line_bot, HAIR, 0.5, dashed=True)
    # steps
    for i, step in enumerate(steps):
        if len(step) == 3:
            f, t, lbl = step; col = NAVY
        else:
            f, t, lbl, col = step
        y = line_top + step_h * (i + 0.5)
        note_h = min(step_h - Inches(0.04), Inches(0.26))
        if f == t:
            # self note (lane 폭 안에 맞춤)
            cx = x0 + lane_w * f + lane_w / 2
            w = lane_w - Inches(0.2)
            rrect(slide, cx - w / 2, y - note_h / 2, w, note_h,
                  LIGHT, line=col, r=0.2)
            txt(slide, cx - w / 2, y - note_h / 2, w, note_h,
                lbl, size=7, color=INK,
                align=PP_ALIGN.CENTER, vcenter=True)
        else:
            x1 = x0 + lane_w * f + lane_w / 2
            x2 = x0 + lane_w * t + lane_w / 2
            arrow(slide, x1, y, x2, y, col, weight=1.25)
            mid = (x1 + x2) / 2
            label_w = Inches(3.0)
            txt(slide, mid - label_w / 2, y - step_h * 0.55,
                label_w, step_h * 0.5,
                lbl, size=7, color=col,
                align=PP_ALIGN.CENTER, bold=True, vcenter=True)


# ─────────────────────────────────────────────
# 5. Sequence: 약속 생성
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 4, "Sequence — 그룹 약속 생성",
       "Mobile → FastAPI → Postgres + 트리거 fanout")
lanes = [
    ("사용자", MUTED),
    ("Mobile", BRAND),
    ("FastAPI", SUCCESS),
    ("Postgres", HERMES),
    ("Realtime", PURPLE),
]
steps = [
    (0, 1, "+ 새 약속"),
    (1, 1, "MeetupCreate 화면"),
    (0, 1, "그룹 선택"),
    (1, 2, "GET /api/groups/{id}/members"),
    (2, 3, "SELECT (RLS)"),
    (3, 2, "members"),
    (2, 1, "200 + members"),
    (0, 1, "장소·멤버·저장"),
    (1, 2, "POST /api/meetups"),
    (2, 3, "BEGIN; INSERT meetups"),
    (3, 3, "trigger: participants/chat_room/reminder INSERT"),
    (3, 3, "INSERT 추가 participants"),
    (3, 2, "COMMIT"),
    (2, 1, "201 + meetup"),
    (3, 4, "NOTIFY"),
    (1, 0, "약속 상세 화면"),
]
seq_diagram(s, lanes, steps)
footer(s, 5, TOTAL)

# ─────────────────────────────────────────────
# 6. Sequence: 위치 공유
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 5, "Sequence — 위치 공유 (백그라운드)",
       "starts_at - share_minutes_before 자동 시작 → ends_at 자동 종료")
lanes = [
    ("Mobile (BG)", BRAND),
    ("Supabase", HERMES),
    ("다른 멤버 Mobile", PURPLE),
    ("pg_cron", SUCCESS),
]
steps = [
    (0, 0, "Local timer fires"),
    (0, 0, "expo-location BG task 시작"),
    (0, 1, "INSERT location_pings (10s)"),
    (1, 1, "RLS: own user_id only"),
    (1, 2, "Realtime publish"),
    (2, 2, "지도 핀 업데이트"),
    (0, 1, "...반복..."),
    (0, 0, "ends_at 도달 → BG task 정지"),
    (3, 1, "tick_meetup_status: status='ended'"),
    (3, 1, "chat_rooms.archived_at = now()"),
    (3, 1, "cleanup_pings (24h 후)"),
]
seq_diagram(s, lanes, steps)
footer(s, 6, TOTAL)

# ─────────────────────────────────────────────
# 7. Sequence: 채팅 메시지
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 6, "Sequence — 채팅 (텍스트/이미지/장소)",
       "송신은 FastAPI 경유, 수신은 Supabase Realtime 직결")
lanes = [
    ("사용자", MUTED),
    ("Mobile", BRAND),
    ("FastAPI", SUCCESS),
    ("Postgres", HERMES),
    ("다른 멤버", PURPLE),
]
steps = [
    (0, 1, "+ 장소 선택"),
    (1, 1, "PlacePicker → place 선택"),
    (1, 2, "POST /api/chat/{room}/messages"),
    (2, 2, "권한·archived 체크"),
    (2, 3, "INSERT messages"),
    (3, 2, "ok"),
    (2, 1, "201"),
    (1, 1, "optimistic add"),
    (3, 4, "Realtime NOTIFY"),
    (4, 4, "메시지 수신 + 푸시 알림"),
]
seq_diagram(s, lanes, steps)
footer(s, 7, TOTAL)

# ─────────────────────────────────────────────
# 8. Sequence: 초대 수락
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 7, "Sequence — 초대 수락 (그룹)",
       "딥링크 → 메타 조회 → 수락 RPC")
lanes = [
    ("사용자", MUTED),
    ("Mobile", BRAND),
    ("FastAPI", SUCCESS),
    ("Postgres", HERMES),
]
steps = [
    (0, 1, "deep link 탭"),
    (1, 1, "/invite/abc123 InviteAccept"),
    (1, 2, "GET /api/invites/abc123"),
    (2, 3, "SELECT invites"),
    (3, 3, "만료/소진/취소 체크"),
    (3, 2, "메타 반환"),
    (2, 1, "200 + 메타"),
    (0, 1, "수락"),
    (1, 2, "POST /api/invites/abc123/accept"),
    (2, 3, "BEGIN; SELECT FOR UPDATE invites"),
    (3, 3, "INSERT group_members"),
    (3, 3, "UPDATE used_count++"),
    (3, 2, "COMMIT"),
    (2, 1, "200"),
    (1, 0, "GroupDetail로 이동"),
]
seq_diagram(s, lanes, steps)
footer(s, 8, TOTAL)

# ─────────────────────────────────────────────
# 9. Meetup state transition
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 8, "State Transition — Meetup",
       "scheduled → active → ended (또는 cancelled)")

states = [
    (Inches(2.0), Inches(3.0), "scheduled", BRAND, "기본 상태"),
    (Inches(6.5), Inches(3.0), "active", SUCCESS,
     "starts_at <= now()"),
    (Inches(11.0), Inches(3.0), "ended", MUTED,
     "ends_at <= now()"),
    (Inches(6.5), Inches(5.5), "cancelled", DANGER,
     "creator/admin 취소"),
]
for x, y, name, c, desc in states:
    rrect(s, x, y, Inches(2.2), Inches(1.4), WHITE, line=c, r=0.3)
    rect(s, x, y, Inches(2.2), Inches(0.45), c)
    txt(s, x, y + Inches(0.10), Inches(2.2), Inches(0.3),
        name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), y + Inches(0.6),
        Inches(1.9), Inches(0.7),
        desc, size=10, color=INK_2, align=PP_ALIGN.CENTER)

# arrows
arrow(s, Inches(4.2), Inches(3.7), Inches(6.5), Inches(3.7), NAVY)
txt(s, Inches(4.4), Inches(3.4), Inches(2.0), Inches(0.3),
    "pg_cron tick", size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

arrow(s, Inches(8.7), Inches(3.7), Inches(11.0), Inches(3.7), NAVY)
txt(s, Inches(8.9), Inches(3.4), Inches(2.0), Inches(0.3),
    "pg_cron tick", size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

arrow(s, Inches(3.1), Inches(4.4), Inches(6.5), Inches(5.5), DANGER)
txt(s, Inches(3.5), Inches(4.7), Inches(2.5), Inches(0.3),
    "취소", size=9, color=DANGER, bold=True)

arrow(s, Inches(7.6), Inches(4.4), Inches(7.6), Inches(5.5), DANGER)
txt(s, Inches(7.8), Inches(4.7), Inches(2.5), Inches(0.3),
    "취소", size=9, color=DANGER, bold=True)

# side effects
rrect(s, Inches(0.5), Inches(6.30),
      Inches(12.3), Inches(0.70), LIGHT, r=0.15)
txt(s, Inches(0.7), Inches(6.36),
    Inches(12), Inches(0.3),
    "부수 효과", size=10, bold=True, color=NAVY)
txt(s, Inches(0.7), Inches(6.62),
    Inches(12), Inches(0.32),
    "active → ended: chat_room.archived_at = now()    "
    "·    * → cancelled: archived + 참여자 푸시    ·    "
    "scheduled → active: 별도 작업 없음 (모바일 자체 트리거)",
    size=9, color=INK_2)
footer(s, 9, TOTAL)

# ─────────────────────────────────────────────
# 10. Backend modules
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 9, "Backend Modules",
       "MeetPod/backend/ — PickPod와 동일 컨벤션")

cols = [
    ("routers/", BRAND, [
        "auth.py", "profiles.py", "friendships.py",
        "invites.py", "groups.py", "meetups.py",
        "reminders.py", "chat.py",
    ]),
    ("services/", SUCCESS, [
        "invite_service",
        "group_service",
        "meetup_service",
        "chat_service",
        "reminder_service",
        "push_service",
        "place_service",
    ]),
    ("dependencies/", HERMES, [
        "auth.py",
        "  (Bearer JWT → user_id)",
        "permissions.py",
        "  (require_owner_or_admin)",
    ]),
    ("utils/", PURPLE, [
        "jwt_utils.py",
        "supabase_client.py",
        "db.py",
        "  (single 헬퍼)",
        "invite_code.py",
        "  (Crockford 8자)",
    ]),
    ("models/", MUTED, [
        "auth.py",
        "profile.py",
        "group.py",
        "meetup.py",
        "chat.py",
        "invite.py",
    ]),
]
x0 = Inches(0.5)
cw = Inches(2.5)
gap = Inches(0.05)
y0 = Inches(1.3)
for i, (name, color, items) in enumerate(cols):
    x = x0 + (cw + gap) * i
    rect(s, x, y0, cw, Inches(0.55), color)
    txt(s, x, y0 + Inches(0.13), cw, Inches(0.3),
        name, size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    rect(s, x, y0 + Inches(0.55), cw, Inches(4.8), SURFACE)
    for j, it in enumerate(items):
        txt(s, x + Inches(0.2), y0 + Inches(0.7 + j * 0.45),
            cw - Inches(0.4), Inches(0.4),
            "• " + it if not it.startswith("  ") else it,
            size=10, color=INK if not it.startswith("  ") else MUTED)

# entry note
rrect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55),
      LIGHT, line=NAVY, r=0.15)
txt(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
    "Vercel entry: api/index.py → app/main.py (FastAPI). "
    "라우터 thin · 비즈니스 로직 services/. 모든 라우트 prefix /api.",
    size=10, color=INK_2)
footer(s, 10, TOTAL)

# ─────────────────────────────────────────────
# 11. Mobile modules
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 10, "Mobile Modules", "MeetPod/mobile/src/")

mcols = [
    ("api/", BRAND, [
        "client.ts (axios)",
        "  baseURL · JWT interceptor",
        "  401 → re-login",
        "auth · profiles · groups",
        "meetups · invites · chat",
        "reminders",
    ]),
    ("lib/", HERMES, [
        "supabase.ts",
        "  Supabase JS singleton",
        "location_tracker.ts",
        "  expo-location TaskManager",
        "push_registrar.ts",
        "deep_link.ts (Linking)",
        "time.ts (day.js)",
    ]),
    ("store/", SUCCESS, [
        "auth_store",
        "  session, profile",
        "meetups_store",
        "  약속 목록 캐시",
        "chat_store",
        "  room별 + Realtime lifecycle",
        "friends_store",
    ]),
    ("navigation/", PURPLE, [
        "root.tsx",
        "  Auth or MainTab",
        "auth_stack.tsx",
        "main_tab.tsx",
        "  Meetups · Groups",
        "  Chats · Me",
    ]),
    ("screens/components/theme/", MUTED, [
        "screens/auth · meetups",
        "screens/groups · chats · me",
        "components/",
        "  Avatar · Pill · MapPin",
        "  MessageBubble · PlaceCard",
        "theme/ (design tokens)",
    ]),
]
x0 = Inches(0.5)
cw = Inches(2.5)
gap = Inches(0.05)
y0 = Inches(1.3)
for i, (name, color, items) in enumerate(mcols):
    x = x0 + (cw + gap) * i
    rect(s, x, y0, cw, Inches(0.55), color)
    txt(s, x, y0 + Inches(0.13), cw, Inches(0.3),
        name, size=11, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    rect(s, x, y0 + Inches(0.55), cw, Inches(5.0), SURFACE)
    for j, it in enumerate(items):
        txt(s, x + Inches(0.2), y0 + Inches(0.7 + j * 0.42),
            cw - Inches(0.4), Inches(0.4),
            "• " + it if not it.startswith("  ") else it,
            size=9, color=INK if not it.startswith("  ") else MUTED)

footer(s, 11, TOTAL)

# ─────────────────────────────────────────────
# 12. API table 1
# ─────────────────────────────────────────────
def api_slide(page, num, title, subtitle, rows):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    header(s, num, title, subtitle)
    hy = Inches(1.2)
    rect(s, Inches(0.5), hy, Inches(12.3), Inches(0.42), NAVY)
    headers = [("Method", 1.2), ("Path", 5.5), ("설명", 3.6), ("권한", 2.0)]
    cx = Inches(0.5)
    for h, w in headers:
        txt(s, cx + Inches(0.15), hy + Inches(0.10),
            Inches(w), Inches(0.3),
            h, size=10, bold=True, color=WHITE)
        cx += Inches(w)
    method_color = {
        "GET": ACCENT, "POST": SUCCESS, "PATCH": WARN,
        "DELETE": DANGER, "PUT": PURPLE,
    }
    for i, (m, p, d, r) in enumerate(rows):
        y = hy + Inches(0.50) + Inches(i * 0.34)
        bg = LIGHT if i % 2 == 0 else WHITE
        rect(s, Inches(0.5), y, Inches(12.3), Inches(0.32), bg)
        cx = Inches(0.5)
        # method pill
        rrect(s, cx + Inches(0.1), y + Inches(0.06),
              Inches(0.85), Inches(0.20),
              method_color.get(m, NAVY), r=0.5)
        txt(s, cx + Inches(0.1), y + Inches(0.07),
            Inches(0.85), Inches(0.2),
            m, size=8, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        cx += Inches(1.2)
        txt(s, cx + Inches(0.15), y + Inches(0.08),
            Inches(5.5), Inches(0.3),
            p, size=9, color=INK)
        cx += Inches(5.5)
        txt(s, cx + Inches(0.15), y + Inches(0.08),
            Inches(3.6), Inches(0.3),
            d, size=9, color=INK_2)
        cx += Inches(3.6)
        txt(s, cx + Inches(0.15), y + Inches(0.08),
            Inches(2.0), Inches(0.3),
            r, size=9, color=MUTED)
    footer(s, page, TOTAL)


api_rows_1 = [
    ("POST", "/api/auth/bootstrap", "profile 행 보장", "self"),
    ("PATCH", "/api/auth/handle", "핸들 1회 설정", "self"),
    ("GET", "/api/profiles/me", "내 프로필", "self"),
    ("PATCH", "/api/profiles/me", "display_name/avatar/token", "self"),
    ("GET", "/api/friends", "친구 목록", "self"),
    ("DELETE", "/api/friends/{user_id}", "친구 끊기", "self"),
    ("POST", "/api/invites", "초대 코드 발급", "self / admin"),
    ("GET", "/api/invites/{code}", "코드 메타", "anon"),
    ("POST", "/api/invites/{code}/accept", "수락", "self"),
    ("DELETE", "/api/invites/{code}", "발급자 취소", "inviter"),
    ("GET", "/api/groups", "내 그룹 목록", "self"),
    ("POST", "/api/groups", "그룹 생성", "self"),
    ("GET", "/api/groups/{id}", "그룹 상세", "member"),
    ("PATCH", "/api/groups/{id}", "그룹 수정", "owner/admin"),
    ("DELETE", "/api/groups/{id}", "그룹 해체", "owner"),
    ("GET", "/api/groups/{id}/members", "멤버 목록", "member"),
]
api_slide(12, 11, "API Endpoints (1/2)",
          "Auth · Profile · Friend · Invite · Group", api_rows_1)

api_rows_2 = [
    ("PATCH", "/api/groups/{id}/members/{uid}/role",
     "admin↔member", "owner"),
    ("DELETE", "/api/groups/{id}/members/{uid}", "추방", "owner/admin"),
    ("POST", "/api/groups/{id}/leave", "본인 탈퇴", "self"),
    ("POST", "/api/groups/{id}/transfer-owner", "owner 위임", "owner"),
    ("GET", "/api/meetups", "내 약속 (다가오는/오늘/지난)", "self"),
    ("POST", "/api/meetups", "약속 생성", "self"),
    ("GET", "/api/meetups/{id}", "약속 상세", "participant"),
    ("PATCH", "/api/meetups/{id}", "정보 수정", "creator/admin"),
    ("POST", "/api/meetups/{id}/cancel", "취소", "creator/admin"),
    ("POST", "/api/meetups/{id}/participants", "멤버 추가",
     "creator/admin"),
    ("DELETE", "/api/meetups/{id}/participants/{uid}",
     "제거/탈퇴", "self/admin"),
    ("PATCH", "/api/meetups/{id}/share-location",
     "본인 위치공유 토글", "self"),
    ("PUT", "/api/meetups/{id}/reminders/me", "본인 알림 시각", "self"),
    ("POST", "/api/chat/{room_id}/messages",
     "메시지 전송", "room member"),
    ("PATCH", "/api/chat/messages/{id}", "본인 메시지 편집", "sender"),
    ("DELETE", "/api/chat/messages/{id}", "soft delete", "sender"),
]
api_slide(13, 12, "API Endpoints (2/2)",
          "Group(cont.) · Meetup · Reminder · Chat", api_rows_2)


# ─────────────────────────────────────────────
# 14. Business rules
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 13, "Business Rules", "강제 위치 정리 — DB · Trigger · API 분담")

rules = [
    ("핸들 unique + 정규식", "DB UNIQUE + CHECK + API 검증", BRAND),
    ("그룹은 owner 1명 필수", "trigger (마지막 owner 탈퇴 차단)", SUCCESS),
    ("1회성 약속은 그룹 없음", "meetups.group_id NULL 허용 + API 분기", BRAND),
    ("약속 종료는 ends_at + pg_cron",
     "트리거 X — 클라이언트 시간 조작 무관", SUCCESS),
    ("위치 핑은 active 상태에서만", "RLS + meetups.status 체크", HERMES),
    ("채팅은 archived 후도 송수신 가능",
     "API의 archived 체크 제거 (시나리오 11)", BRAND),
    ("푸시 중복 방지",
     "reminder 행 자체가 큐 → 발송 후 DELETE (idempotent)", PURPLE),
    ("초대 코드 충돌 방지",
     "DB UNIQUE + INSERT 충돌 시 재생성 (8자 base32 ≈ 10억)",
     SUCCESS),
]
y0 = Inches(1.2)
rh = Inches(0.6)
for i, (k, v, color) in enumerate(rules):
    y = y0 + Inches(i * 0.66)
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.5), y, Inches(12.3), rh, bg)
    rect(s, Inches(0.5), y, Inches(0.10), rh, color)
    txt(s, Inches(0.7), y + Inches(0.08), Inches(5.5), Inches(0.3),
        k, size=11, bold=True, color=INK)
    txt(s, Inches(6.3), y + Inches(0.18), Inches(6.4), Inches(0.3),
        v, size=10, color=INK_2)
footer(s, 14, TOTAL)


# ─────────────────────────────────────────────
# 15. Non-functional + 16. Security (combined? separate)
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 14, "Non-Functional Requirements", "성능 · 보존 · 가용성")

nfr = [
    ("응답 시간", "p95 API 300ms", "Vercel cold 제외, 단순 SELECT는 직결"),
    ("위치 핑 latency", "end-to-end ≤ 2s", "Supabase Realtime"),
    ("푸시 정확도", "reminder dispatch ±1분", "pg_cron 1분 주기"),
    ("동시 접속", "MVP 1,000명", "Supabase Free/Pro 충분"),
    ("가용성", "99.5%", "Vercel + Supabase SLA"),
    ("location_pings 보존", "24시간", "프라이버시 + 비용"),
    ("messages 보존", "무기한", "soft delete만"),
    ("탈퇴 시 profile", "즉시 익명화", "GDPR-friendly"),
]
hy = Inches(1.2)
rect(s, Inches(0.5), hy, Inches(12.3), Inches(0.42), NAVY)
hdrs = [("영역", 3.5), ("목표", 3.0), ("비고", 5.8)]
cx = Inches(0.5)
for h, w in hdrs:
    txt(s, cx + Inches(0.15), hy + Inches(0.10),
        Inches(w), Inches(0.3),
        h, size=11, bold=True, color=WHITE)
    cx += Inches(w)
for i, (k, v, n) in enumerate(nfr):
    y = hy + Inches(0.50) + Inches(i * 0.55)
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5), bg)
    rect(s, Inches(0.5), y, Inches(0.10), Inches(0.5), HERMES)
    txt(s, Inches(0.7), y + Inches(0.13), Inches(3.3), Inches(0.3),
        k, size=11, bold=True, color=INK)
    txt(s, Inches(4.0), y + Inches(0.13), Inches(2.8), Inches(0.3),
        v, size=11, color=BRAND, bold=True)
    txt(s, Inches(7.0), y + Inches(0.15), Inches(5.5), Inches(0.3),
        n, size=10, color=MUTED)
footer(s, 15, TOTAL)


# ─────────────────────────────────────────────
# 16. Security
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 15, "Security — Threat Model", "위협별 대응")

threats = [
    ("타 사용자 위치 노출",
     "RLS (같은 meetup 참여자만) + share_location 본인 토글", DANGER),
    ("초대 코드 brute force",
     "8자 base32 (10억 공간) + invite 조회 rate limit", DANGER),
    ("JWT 탈취",
     "Refresh token rotation + 클라이언트 secure storage", WARN),
    ("RLS 우회 (service_role 키 노출)",
     "키는 backend env만, 클라이언트는 anon key", DANGER),
    ("채팅 첨부 악용",
     "Storage 정책 + 이미지 크기/타입 검증", WARN),
    ("푸시 스팸",
     "사용자 알림 설정, 약속당 변경/취소 알림 N회 rate guard",
     SUCCESS),
    ("핸들 squatting",
     "14일 미사용 회수 정책 (Phase 2). MVP는 영구", MUTED),
]
y0 = Inches(1.2)
for i, (k, v, color) in enumerate(threats):
    y = y0 + Inches(i * 0.75)
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.65), bg)
    rect(s, Inches(0.5), y, Inches(0.10), Inches(0.65), color)
    txt(s, Inches(0.7), y + Inches(0.08), Inches(0.5), Inches(0.5),
        "⚠", size=18, color=color)
    txt(s, Inches(1.3), y + Inches(0.06), Inches(11.4), Inches(0.3),
        k, size=12, bold=True, color=INK)
    txt(s, Inches(1.3), y + Inches(0.32), Inches(11.4), Inches(0.3),
        v, size=10, color=INK_2)
footer(s, 16, TOTAL)


# ─────────────────────────────────────────────
# 17. Deploy + Observability + Test
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 16, "Deploy · Observability · Test",
       "환경 분리 · 관측 · 테스트 전략")

# Environments
rrect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(2.7),
      SURFACE, line=BRAND, r=0.1)
txt(s, Inches(0.7), Inches(1.30), Inches(5.6), Inches(0.3),
    "Environments", size=13, bold=True, color=BRAND)
envs = [("dev", "meetpod-dev · local uvicorn · Expo Go"),
        ("staging", "meetpod-staging · Vercel Preview · EAS internal"),
        ("prod", "meetpod-prod · Vercel Production · EAS production")]
for i, (k, v) in enumerate(envs):
    y = Inches(1.7 + i * 0.65)
    rrect(s, Inches(0.7), y, Inches(1.2), Inches(0.5),
          BRAND, r=0.5)
    txt(s, Inches(0.7), y + Inches(0.13), Inches(1.2), Inches(0.3),
        k, size=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    txt(s, Inches(2.0), y + Inches(0.15), Inches(4.4), Inches(0.4),
        v, size=10, color=INK_2)

# Observability
rrect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(2.7),
      SURFACE, line=HERMES, r=0.1)
txt(s, Inches(7.0), Inches(1.30), Inches(5.6), Inches(0.3),
    "Observability", size=13, bold=True, color=HERMES)
obs = [
    ("Backend 로그", "Vercel logs"),
    ("DB 쿼리", "Supabase Studio"),
    ("Realtime", "Supabase Dashboard"),
    ("푸시", "Edge Function 로그 + outbox (Phase 2)"),
    ("모바일 크래시", "Sentry (Expo 통합)"),
    ("사용자 분석", "PostHog/Amplitude (Phase 2)"),
]
for i, (k, v) in enumerate(obs):
    y = Inches(1.7 + i * 0.36)
    txt(s, Inches(7.0), y, Inches(2.5), Inches(0.3),
        "• " + k, size=10, bold=True, color=INK)
    txt(s, Inches(9.5), y, Inches(3.2), Inches(0.3),
        v, size=10, color=MUTED)

# Testing
rrect(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7),
      SURFACE, line=SUCCESS, r=0.1)
txt(s, Inches(0.7), Inches(4.20), Inches(12), Inches(0.3),
    "Testing Strategy", size=13, bold=True, color=SUCCESS)
tests = [
    ("단위 (backend)", "pytest", "service 함수, 유효성 검증"),
    ("RLS", "pgTAP / SQL 시나리오", "권한 우회 시도, 정책 정합성"),
    ("통합", "pytest + Supabase local", "라우터 ↔ DB ↔ RLS"),
    ("모바일 단위", "Jest + RNTL", "store, util, 작은 컴포넌트"),
    ("수동 E2E", "TestFlight / EAS internal",
     "백그라운드 위치 (시뮬레이터 부정확)"),
    ("부하", "Locust", "채팅 fan-out, 위치 핑 INSERT"),
]
hy = Inches(4.65)
rect(s, Inches(0.7), hy, Inches(11.9), Inches(0.4), NAVY)
for i, (h, w) in enumerate([("레벨", 3.0), ("도구", 3.5), ("범위", 5.4)]):
    cx = Inches(0.7) + Inches(sum(
        [3.0, 3.5][:i]))
    txt(s, cx + Inches(0.15), hy + Inches(0.08),
        Inches(w), Inches(0.3),
        h, size=10, bold=True, color=WHITE)
for i, (a, b, c) in enumerate(tests):
    y = hy + Inches(0.45 + i * 0.32)
    bg = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.30), bg)
    txt(s, Inches(0.85), y + Inches(0.06), Inches(2.85),
        Inches(0.3), a, size=10, bold=True, color=INK)
    txt(s, Inches(3.85), y + Inches(0.06), Inches(3.35),
        Inches(0.3), b, size=10, color=BRAND, bold=True)
    txt(s, Inches(7.20), y + Inches(0.06), Inches(5.4),
        Inches(0.3), c, size=10, color=MUTED)

footer(s, 17, TOTAL)


# ─────────────────────────────────────────────
# 18. Risks Top 5
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
header(s, 17, "Top Risks", "출시 전 반드시 검증할 5가지")

risks = [
    ("iOS 백그라운드 위치 권한 거부",
     "위치공유 동작 안 함",
     "OS 설정 가이드 + 포그라운드 fallback", DANGER),
    ("Kakao 로그인 Supabase 비호환",
     "가입 마찰",
     "Custom OIDC 또는 백엔드 검증 fallback", WARN),
    ("Google Places API 비용 폭증",
     "운영비",
     "디바운스, 캐싱, 일일 한도 알림", WARN),
    ("위치 핑 INSERT 양 폭증",
     "DB I/O",
     "인덱스, TTL 단축, 핑 주기 증가 옵션", DANGER),
    ("푸시 알림 과다 → 사용자 이탈",
     "리텐션",
     "사용자별 채널 토글, 빈도 가이드", WARN),
]
y0 = Inches(1.2)
for i, (risk, impact, mitig, color) in enumerate(risks):
    y = y0 + Inches(i * 1.05)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(0.95),
          SURFACE, line=color, r=0.1)
    # rank
    oval(s, Inches(0.65), y + Inches(0.20),
         Inches(0.55), Inches(0.55), color)
    txt(s, Inches(0.65), y + Inches(0.30),
        Inches(0.55), Inches(0.4),
        f"#{i + 1}", size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    # risk
    txt(s, Inches(1.4), y + Inches(0.10), Inches(6.0), Inches(0.4),
        risk, size=13, bold=True, color=INK)
    # impact
    rrect(s, Inches(1.4), y + Inches(0.50), Inches(2.5), Inches(0.30),
          color, r=0.5)
    txt(s, Inches(1.4), y + Inches(0.53),
        Inches(2.5), Inches(0.25),
        f"영향: {impact}", size=9, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER)
    # mitigation
    txt(s, Inches(4.2), y + Inches(0.50), Inches(8.5),
        Inches(0.4),
        f"→ {mitig}", size=11, color=INK_2)
footer(s, 18, TOTAL)


out = "d:/Workspace/CPKWorks/MeetPod/docs/MeetPod_시스템설계.pptx"
prs.save(out)
